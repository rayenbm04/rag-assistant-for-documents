import os, shutil, base64, asyncio, hashlib, json, re, uuid, subprocess
import requests
from bs4 import BeautifulSoup
from pptx import Presentation as PptxPresentation
from datetime import datetime, timedelta
import nest_asyncio  # allows AutoMergingRetriever to run inside FastAPI's event loop
import pdfplumber
from PIL import Image
from docx import Document as DocxDocument
from concurrent.futures import ThreadPoolExecutor
from fastapi import FastAPI, UploadFile, File, HTTPException, Depends
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse
from fastapi.security import OAuth2PasswordBearer
from pydantic import BaseModel
from dotenv import load_dotenv
from jose import JWTError, jwt
from passlib.context import CryptContext
from sqlalchemy import create_engine, Column, String, DateTime
from sqlalchemy.orm import declarative_base, sessionmaker
import chromadb
import ollama
from llama_index.core import VectorStoreIndex, Settings, Document, PromptTemplate
from llama_index.core.vector_stores.types import MetadataFilters, MetadataFilter, FilterOperator, FilterCondition
from llama_index.core.retrievers import QueryFusionRetriever
from llama_index.core.schema import TextNode
from llama_index.retrievers.bm25 import BM25Retriever
from llama_index.core.storage.storage_context import StorageContext

# Parent-child retrieval — graceful fallback if not available in this build
try:
    from llama_index.core.retrievers import AutoMergingRetriever
    from llama_index.core.node_parser import HierarchicalNodeParser, get_leaf_nodes
    from llama_index.core.storage.docstore import SimpleDocumentStore
    PARENT_CHILD_AVAILABLE = True
except ImportError as _pc_err:
    print(f"[warn] Parent-child retrieval not available ({_pc_err}) — using flat chunking")
    PARENT_CHILD_AVAILABLE = False
from llama_index.vector_stores.chroma import ChromaVectorStore
from llama_index.llms.ollama import Ollama
from llama_index.embeddings.ollama import OllamaEmbedding

load_dotenv()

LLM_MODEL        = os.getenv("LLM_MODEL",           "qwen2.5:7b")
EMBED_MODEL      = os.getenv("EMBED_MODEL",          "nomic-embed-text")
VISION_MODEL     = os.getenv("VISION_MODEL",         "qwen2.5vl:7b")
OPENAI_API_KEY      = os.getenv("OPENAI_API_KEY",        "")
OPENAI_MODEL        = os.getenv("OPENAI_MODEL",          "gpt-4o")
GROQ_API_KEY        = os.getenv("GROQ_API_KEY",          "")
GROQ_MODEL          = os.getenv("GROQ_MODEL",            "llama-3.3-70b-versatile")
GEMINI_API_KEY      = os.getenv("GEMINI_API_KEY",        "")
GEMINI_VISION_MODEL = os.getenv("GEMINI_VISION_MODEL",   "gemini-1.5-flash")
UPLOAD_DIR       = os.getenv("UPLOAD_DIR",           "./uploads")
CHROMA_DIR       = os.getenv("CHROMA_DIR",           "./chroma_db")
SLIDES_CACHE_DIR = os.getenv("SLIDES_CACHE_DIR",     "./slides_cache")
ALLOWED_ORIGINS  = os.getenv("ALLOWED_ORIGINS",      "http://localhost:5173").split(",")
SIMILARITY_TOP_K = int(os.getenv("SIMILARITY_TOP_K", "4"))
MAX_UPLOAD_MB    = int(os.getenv("MAX_UPLOAD_SIZE_MB", "50"))
OLLAMA_BASE_URL  = os.getenv("OLLAMA_BASE_URL",      "http://localhost:11434")
ENABLE_EVAL        = os.getenv("ENABLE_EVAL",          "true").lower() == "true"
ENABLE_HYDE        = os.getenv("ENABLE_HYDE",          "true").lower() == "true"
ENABLE_MULTI_QUERY = os.getenv("ENABLE_MULTI_QUERY",   "true").lower() == "true"
MULTI_QUERY_N      = int(os.getenv("MULTI_QUERY_N",    "3"))
ENABLE_RERANK      = os.getenv("ENABLE_RERANK",        "true").lower() == "true"
RERANK_MODEL       = os.getenv("RERANK_MODEL",         "cross-encoder/ms-marco-MiniLM-L-6-v2")
PARENT_CHUNK_SIZE  = int(os.getenv("PARENT_CHUNK_SIZE", "512"))
CHILD_CHUNK_SIZE   = int(os.getenv("CHILD_CHUNK_SIZE",  "128"))
NODE_STORE_DIR     = os.getenv("NODE_STORE_DIR",        "./node_store")
SECRET_KEY             = os.getenv("SECRET_KEY",             "change-me-in-production")
if SECRET_KEY == "change-me-in-production":
    import warnings
    warnings.warn(
        "SECRET_KEY is set to the default insecure value. "
        "Set a random SECRET_KEY in .env before deploying.",
        stacklevel=1,
    )
ACCESS_TOKEN_EXPIRE_DAYS = int(os.getenv("ACCESS_TOKEN_EXPIRE_DAYS", "7"))
DATABASE_URL           = os.getenv("DATABASE_URL",           "sqlite:///./rag_users.db")

_engine      = create_engine(DATABASE_URL, connect_args={"check_same_thread": False})
_SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=_engine)
_Base        = declarative_base()

class UserModel(_Base):
    __tablename__ = "users"
    id              = Column(String, primary_key=True, default=lambda: str(uuid.uuid4()))
    email           = Column(String, unique=True, index=True, nullable=False)
    firstname       = Column(String, nullable=False)
    lastname        = Column(String, nullable=False)
    hashed_password = Column(String, nullable=False)
    role            = Column(String, default="user")
    created_at      = Column(DateTime, default=datetime.utcnow)

_Base.metadata.create_all(bind=_engine)

_pwd_ctx      = CryptContext(schemes=["bcrypt"], deprecated="auto")
_oauth2_scheme = OAuth2PasswordBearer(tokenUrl="/auth/login")

def _hash_pw(pw: str) -> str:           return _pwd_ctx.hash(pw)
def _verify_pw(plain: str, h: str) -> bool: return _pwd_ctx.verify(plain, h)

def _make_token(user: UserModel) -> str:
    expire = datetime.utcnow() + timedelta(days=ACCESS_TOKEN_EXPIRE_DAYS)
    return jwt.encode(
        {"sub": user.id, "email": user.email, "role": user.role, "exp": expire},
        SECRET_KEY, algorithm="HS256",
    )

def _db_session():
    db = _SessionLocal()
    try:
        yield db
    finally:
        db.close()

async def get_current_user(token: str = Depends(_oauth2_scheme)):
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=["HS256"])
        user_id: str = payload.get("sub")
        if not user_id:
            raise HTTPException(401, "Invalid token")
    except JWTError:
        raise HTTPException(401, "Invalid token")
    db = _SessionLocal()
    try:
        user = db.query(UserModel).filter(UserModel.id == user_id).first()
    finally:
        db.close()
    if not user:
        raise HTTPException(401, "User not found")
    return user

def _require_admin(current_user: UserModel = Depends(get_current_user)):
    if current_user.role != "admin":
        raise HTTPException(403, "Admin access required")
    return current_user

_FILE_OWNERS_PATH = "./file_owners.json"

def _load_file_owners() -> dict:
    if os.path.exists(_FILE_OWNERS_PATH):
        with open(_FILE_OWNERS_PATH) as f:
            return json.load(f)
    return {}

def _save_file_owners():
    with open(_FILE_OWNERS_PATH, "w") as f:
        json.dump(file_owners, f)

file_owners: dict[str, str] = _load_file_owners()   # {filename: user_id}

class RegisterRequest(BaseModel):
    email: str
    password: str
    firstname: str
    lastname: str

class LoginRequest(BaseModel):
    email: str
    password: str

class UrlIngestRequest(BaseModel):
    url: str


class HistoryEntry(BaseModel):
    question: str
    answer: str

class Question(BaseModel):
    question: str
    history: list[HistoryEntry] = []
    files: list[str] = []
    provider: str = "local"       # "local" (Ollama) or "cloud"
    groq_model: str | None = None # override Groq model for this request only
    fast: bool = False            # skip HyDE/multi-query/condense/eval for batch evals

cancelled_files = set()   # filenames that should stop indexing

_llm_extra = {"num_gpu": 99}
if LLM_MODEL.lower().startswith("qwen3"):
    _llm_extra["think"] = False
Settings.llm = Ollama(model=LLM_MODEL, base_url=OLLAMA_BASE_URL, request_timeout=120.0, additional_kwargs=_llm_extra)
Settings.embed_model = OllamaEmbedding(model_name=EMBED_MODEL, base_url=OLLAMA_BASE_URL, ollama_additional_kwargs={"num_gpu": 99})

try:
    from llama_index.llms.openai import OpenAI as OpenAILLM
    _OPENAI_AVAILABLE = True
except ImportError:
    _OPENAI_AVAILABLE = False

_PROVIDER_FILE    = os.path.join(os.getenv("UPLOAD_DIR", "./uploads"), ".provider")
_GROQ_USAGE_FILE  = os.path.join(os.getenv("UPLOAD_DIR", "./uploads"), ".groq_usage.json")

def _load_provider() -> str:
    try:
        return open(_PROVIDER_FILE).read().strip()
    except Exception:
        return "local"

def _save_provider(p: str):
    try:
        os.makedirs(os.path.dirname(_PROVIDER_FILE) or ".", exist_ok=True)
        open(_PROVIDER_FILE, "w").write(p)
    except Exception:
        pass

_active_provider: str = _load_provider()
_last_groq_model: str = GROQ_MODEL  # updated per-request, shown in dashboard

class _CompatLLM:
    """Thin async wrapper around any OpenAI-compatible API.
    Bypasses LlamaIndex model-name validation and returns objects
    compatible with our acomplete/astream_complete call sites."""

    def __init__(self, api_key: str, model: str, base_url: str = "https://api.openai.com/v1"):
        from openai import AsyncOpenAI
        self._client   = AsyncOpenAI(api_key=api_key, base_url=base_url)
        self._model    = model
        self._is_groq  = "groq.com" in base_url

    def _track_groq(self, usage, headers=None):
        """Update per-model groq_token_usage from response headers (authoritative)
        with our own cumulative count as fallback. Persists to disk."""
        if not self._is_groq:
            return
        import datetime
        today = datetime.date.today().isoformat()
        if groq_token_usage.get("date") != today:
            groq_token_usage["date"]   = today
            groq_token_usage["models"] = {}
        m = groq_token_usage["models"].setdefault(
            self._model, {"prompt": 0, "completion": 0, "used_actual": None, "limit_actual": None}
        )
        # Our own cumulative count (fallback)
        if usage:
            m["prompt"]     += getattr(usage, "prompt_tokens",     0)
            m["completion"] += getattr(usage, "completion_tokens", 0)
        # Authoritative values from Groq response headers
        if headers:
            try:
                limit     = headers.get("x-ratelimit-limit-tokens-day")
                remaining = headers.get("x-ratelimit-remaining-tokens-day")
                used      = headers.get("x-ratelimit-used-tokens-day")
                if limit:
                    m["limit_actual"] = int(limit)
                if used:
                    m["used_actual"] = int(used)
                elif limit and remaining:
                    m["used_actual"] = int(limit) - int(remaining)
                # Per-minute (TPM) stats — track separately so UI can warn about TPM limits
                tpm_limit     = headers.get("x-ratelimit-limit-tokens")
                tpm_remaining = headers.get("x-ratelimit-remaining-tokens")
                tpm_reset     = headers.get("x-ratelimit-reset-tokens")
                if tpm_limit is not None:
                    m["tpm_limit"]     = int(tpm_limit)
                if tpm_remaining is not None:
                    m["tpm_remaining"] = int(tpm_remaining)
                if tpm_reset is not None:
                    m["tpm_reset"]     = tpm_reset  # e.g. "1.234s"
            except (ValueError, TypeError):
                pass
        _save_groq_usage()
        # Reset the daily-limit flag if the date changed (new day = new quota)
        global _groq_daily_limit_hit
        _groq_daily_limit_hit = False

    def _handle_rate_limit(self, e) -> bool:
        """Update token bar from 429 error response headers (works even on failed calls),
        then return True if it's a daily (TPD) limit that won't recover by retrying."""
        # Always try to capture usage from error response headers
        try:
            resp_headers = getattr(getattr(e, "response", None), "headers", None)
            if resp_headers:
                self._track_groq(None, resp_headers)
        except Exception:
            pass
        msg = str(getattr(e, "message", "") or e).lower()
        return "per day" in msg or "tpd" in msg or "tokens per day" in msg

    async def acomplete(self, prompt: str):
        from openai import RateLimitError as _RLE
        for attempt in range(6):
            try:
                raw_resp = await self._client.chat.completions.with_raw_response.create(
                    model=self._model,
                    messages=[{"role": "user", "content": prompt}],
                )
                resp  = raw_resp.parse()
                text  = resp.choices[0].message.content or ""
                usage = resp.usage
                self._track_groq(usage, raw_resp.headers)
                raw = {
                    "prompt_eval_count":  getattr(usage, "prompt_tokens",    0),
                    "eval_count":         getattr(usage, "completion_tokens", 0),
                }
                return type("CR", (), {"text": text, "raw": raw, "__str__": lambda s: text})()
            except _RLE as e:
                is_daily = self._handle_rate_limit(e)
                if is_daily:
                    global _groq_daily_limit_hit
                    _groq_daily_limit_hit = True
                    print(f"  [llm] daily token limit hit — failing immediately")
                    raise RuntimeError("Rate limit reached — Groq daily token quota exhausted. Please try again tomorrow or switch to a different model.")
                wait = 20 * (attempt + 1)
                print(f"  [llm] rate limit (TPM/RPM) — waiting {wait}s (attempt {attempt+1}/6)")
                await asyncio.sleep(wait)
        raise RuntimeError("Rate limit reached — too many requests to Groq. Please wait 1–2 minutes and try again.")

    async def astream_complete(self, prompt: str):
        from openai import RateLimitError as _RLE
        for attempt in range(6):
            try:
                raw_resp = await self._client.chat.completions.with_raw_response.create(
                    model=self._model,
                    messages=[{"role": "user", "content": prompt}],
                    stream=False,
                )
                resp  = raw_resp.parse()
                text  = resp.choices[0].message.content or ""
                self._track_groq(resp.usage, raw_resp.headers)

                async def _gen(t=text):
                    for ch in t:
                        yield type("Chunk", (), {"delta": ch})()

                return _gen()
            except _RLE as e:
                is_daily = self._handle_rate_limit(e)
                if is_daily:
                    global _groq_daily_limit_hit
                    _groq_daily_limit_hit = True
                    print(f"  [llm] daily token limit hit — failing immediately")
                    raise RuntimeError("Rate limit reached — Groq daily token quota exhausted. Please try again tomorrow or switch to a different model.")
                wait = 20 * (attempt + 1)
                print(f"  [llm] rate limit (TPM/RPM) — waiting {wait}s (attempt {attempt+1}/6)")
                await asyncio.sleep(wait)
        raise RuntimeError("Rate limit reached — too many requests to Groq. Please wait 1–2 minutes and try again.")


def _get_llm(provider: str = "local", groq_model: str | None = None):
    """Return the LLM for the given provider. Embeddings always stay local."""
    if provider in ("cloud", "groq"):
        if not GROQ_API_KEY:
            raise HTTPException(status_code=400, detail="GROQ_API_KEY not set in .env")
        model = groq_model or GROQ_MODEL
        return _CompatLLM(api_key=GROQ_API_KEY, model=model, base_url="https://api.groq.com/openai/v1")
    if provider == "openai":
        if not OPENAI_API_KEY:
            raise HTTPException(status_code=400, detail="OPENAI_API_KEY not set in .env")
        return _CompatLLM(api_key=OPENAI_API_KEY, model=OPENAI_MODEL)
    return Settings.llm

app = FastAPI(title="RAG Assistant API")
app.add_middleware(CORSMiddleware,
    allow_origins=ALLOWED_ORIGINS,
    allow_methods=["*"], allow_headers=["*"])

os.makedirs(UPLOAD_DIR,      exist_ok=True)
os.makedirs(NODE_STORE_DIR,  exist_ok=True)
os.makedirs(SLIDES_CACHE_DIR, exist_ok=True)


def _safe_upload_path(filename: str) -> str:
    """Resolve filename inside UPLOAD_DIR and reject any path traversal attempt."""
    safe = os.path.realpath(os.path.join(UPLOAD_DIR, os.path.basename(filename)))
    upload_root = os.path.realpath(UPLOAD_DIR)
    if not safe.startswith(upload_root + os.sep) and safe != upload_root:
        raise HTTPException(400, "Invalid filename")
    return safe


@app.post("/auth/register")
def auth_register(req: RegisterRequest):
    db = _SessionLocal()
    try:
        if db.query(UserModel).filter(UserModel.email == req.email).first():
            raise HTTPException(400, "Email already registered")
        if len(req.password) < 6:
            raise HTTPException(400, "Password must be at least 6 characters")
        is_first = db.query(UserModel).count() == 0
        user = UserModel(
            id=str(uuid.uuid4()),
            email=req.email,
            firstname=req.firstname,
            lastname=req.lastname,
            hashed_password=_hash_pw(req.password),
            role="admin" if is_first else "user",
        )
        db.add(user); db.commit(); db.refresh(user)
    finally:
        db.close()
    token = _make_token(user)
    return {"access_token": token, "token_type": "bearer",
            "user": {"id": user.id, "email": user.email, "firstname": user.firstname, "lastname": user.lastname, "role": user.role}}

@app.post("/auth/login")
def auth_login(req: LoginRequest):
    db = _SessionLocal()
    try:
        user = db.query(UserModel).filter(UserModel.email == req.email).first()
    finally:
        db.close()
    if not user or not _verify_pw(req.password, user.hashed_password):
        raise HTTPException(401, "Invalid email or password")
    token = _make_token(user)
    return {"access_token": token, "token_type": "bearer",
            "user": {"id": user.id, "email": user.email, "firstname": user.firstname, "lastname": user.lastname, "role": user.role}}

@app.get("/auth/me")
def auth_me(current_user: UserModel = Depends(get_current_user)):
    return {"id": current_user.id, "email": current_user.email, "firstname": current_user.firstname, "lastname": current_user.lastname, "role": current_user.role}

@app.get("/auth/users")
def auth_list_users(current_user: UserModel = Depends(_require_admin)):
    db = _SessionLocal()
    try:
        users = db.query(UserModel).all()
        return [{"id": u.id, "email": u.email, "firstname": u.firstname, "lastname": u.lastname, "role": u.role,
                 "created_at": u.created_at.isoformat()} for u in users]
    finally:
        db.close()

@app.patch("/auth/users/{user_id}/role")
def auth_set_role(user_id: str, body: dict,
                  current_user: UserModel = Depends(_require_admin)):
    new_role = body.get("role")
    if new_role not in ("admin", "user"):
        raise HTTPException(400, "role must be 'admin' or 'user'")
    db = _SessionLocal()
    try:
        user = db.query(UserModel).filter(UserModel.id == user_id).first()
        if not user:
            raise HTTPException(404, "User not found")
        user.role = new_role; db.commit()
        return {"id": user.id, "email": user.email, "role": user.role}
    finally:
        db.close()

chroma_client     = chromadb.PersistentClient(path=CHROMA_DIR)
chroma_collection = chroma_client.get_or_create_collection("rag_docs")
vector_store      = ChromaVectorStore(chroma_collection=chroma_collection)

_ds_path = os.path.join(NODE_STORE_DIR, "docstore.json")
if PARENT_CHILD_AVAILABLE:
    if os.path.exists(_ds_path):
        docstore = SimpleDocumentStore.from_persist_path(_ds_path)
        print(f"Loaded docstore: {len(docstore.docs)} nodes")
    else:
        docstore = SimpleDocumentStore()
        print("Created fresh docstore")
else:
    docstore = None

def _persist_docstore():
    if docstore is not None:
        docstore.persist(persist_path=_ds_path)

storage_context = StorageContext.from_defaults(
    vector_store=vector_store,
    **({"docstore": docstore} if docstore is not None else {}),
)

index = None
indexing_status = {}        # {"filename": "indexing" | "ready" | "error"}
indexing_progress = {}      # {"filename": {"current": int, "total": int}}
file_hashes = {}            # {"filename": md5_hex} — used to skip re-indexing unchanged files
executor = ThreadPoolExecutor(max_workers=2)
token_usage = {"prompt": 0, "completion": 0, "requests": 0}
query_stats = {"total": 0, "total_ms": 0}   # questions asked + cumulative response time

# Groq daily token limits per model (tokens-per-day on the free tier)
GROQ_TPD_LIMITS: dict[str, int] = {
    # Active free-tier models (as of 2026-06)
    "llama-3.3-70b-versatile":                             100_000,
    "llama-3.1-8b-instant":                                500_000,
    "meta-llama/llama-4-scout-17b-16e-instruct":           100_000,
    # gemma2-9b-it — decommissioned by Groq, removed
}

# Auxiliary (cheap) model used for condense / HyDE / multi-query / eval scoring
# when the main provider is cloud — preserves 70B quota for the actual answer.
GROQ_AUX_MODEL = "llama-3.1-8b-instant"

def _track_vision_groq_usage(usage, headers=None):
    """Track token usage for the Groq vision model (Llama 4 Scout)."""
    import datetime
    VISION_KEY = "meta-llama/llama-4-scout-17b-16e-instruct"
    today = datetime.date.today().isoformat()
    if groq_token_usage.get("date") != today:
        groq_token_usage["date"]   = today
        groq_token_usage["models"] = {}
    m = groq_token_usage["models"].setdefault(
        VISION_KEY, {"prompt": 0, "completion": 0, "used_actual": None, "limit_actual": None}
    )
    if usage:
        m["prompt"]     += getattr(usage, "prompt_tokens",     0)
        m["completion"] += getattr(usage, "completion_tokens", 0)
    if headers:
        try:
            limit     = headers.get("x-ratelimit-limit-tokens-day")
            remaining = headers.get("x-ratelimit-remaining-tokens-day")
            used      = headers.get("x-ratelimit-used-tokens-day")
            if limit:     m["limit_actual"] = int(limit)
            if used:      m["used_actual"]  = int(used)
            elif limit and remaining: m["used_actual"] = int(limit) - int(remaining)
            tpm_limit     = headers.get("x-ratelimit-limit-tokens")
            tpm_remaining = headers.get("x-ratelimit-remaining-tokens")
            if tpm_limit is not None:     m["tpm_limit"]     = int(tpm_limit)
            if tpm_remaining is not None: m["tpm_remaining"] = int(tpm_remaining)
        except (ValueError, TypeError):
            pass
    _save_groq_usage()

def _load_groq_usage() -> dict:
    """Load today's per-model Groq usage from disk; reset if the date changed."""
    import datetime
    today = datetime.date.today().isoformat()
    try:
        data = json.load(open(_GROQ_USAGE_FILE, encoding="utf-8"))
        if data.get("date") == today:
            return data
    except Exception:
        pass
    return {"date": today, "models": {}}

def _save_groq_usage() -> None:
    try:
        os.makedirs(os.path.dirname(_GROQ_USAGE_FILE) or ".", exist_ok=True)
        json.dump(groq_token_usage, open(_GROQ_USAGE_FILE, "w", encoding="utf-8"))
    except Exception:
        pass

groq_token_usage: dict = _load_groq_usage()   # {"date": "YYYY-MM-DD", "models": {model: {"prompt": n, "completion": n}}}
_groq_daily_limit_hit: bool = False           # Set True when TPD exceeded — skip aux calls to avoid wasting tokens

reranker = None
if ENABLE_RERANK:
    try:
        from sentence_transformers import CrossEncoder
        reranker = CrossEncoder(RERANK_MODEL)
        print(f"Reranker loaded: {RERANK_MODEL}")
    except Exception as _rerank_err:
        print(f"Reranker init failed ({_rerank_err}) — continuing without reranking")


def parse_eval_score(text: str) -> float | None:
    """Extract a 0.0–1.0 score from LLM output. Robust to prose around the number."""
    text = text.strip()
    match = re.search(r'\b(1\.?0*|0?\.\d+|[01])\b', text)
    if match:
        return round(min(max(float(match.group(1)), 0.0), 1.0), 2)
    return None


def record_tokens(result):
    """Extract token counts from a LlamaIndex CompletionResponse and accumulate."""
    raw = getattr(result, "raw", None) or {}
    prompt_tokens     = raw.get("prompt_eval_count", 0) or 0
    completion_tokens = raw.get("eval_count", 0)        or 0
    token_usage["prompt"]     += prompt_tokens
    token_usage["completion"] += completion_tokens
    return prompt_tokens, completion_tokens

if chroma_collection.count() > 0:
    index = VectorStoreIndex.from_vector_store(
        vector_store,
        storage_context=storage_context
    )
    print(f"Loaded index — {chroma_collection.count()} chunks in ChromaDB")



def pil_image_to_base64(pil_image):
    import io
    buffer = io.BytesIO()
    pil_image.save(buffer, format="PNG")
    return base64.b64encode(buffer.getvalue()).decode("utf-8")


def image_to_base64(file_path):
    with open(file_path, "rb") as f:
        return base64.b64encode(f.read()).decode("utf-8")


_UML_KEYWORDS = {
    "uml", "diagram", "diagramme", "schema", "schéma", "architecture",
    "class", "classe", "sequence", "usecase", "use_case", "use-case",
    "activity", "component", "deployment", "statechart", "er_diagram",
    "erd", "flowchart", "dataflow", "dfd",
}

def _is_uml_image(filename: str) -> bool:
    """Return True if the filename suggests a UML or architecture diagram."""
    stem = os.path.splitext(filename.lower())[0]
    # match any keyword as a whole word/token in the stem
    tokens = set(re.split(r'[\s_\-\.]+', stem))
    return bool(tokens & _UML_KEYWORDS)


def _analyze_image_cloud(image_b64: str, prompt: str, max_tokens: int = 3500) -> str:
    """Send image to Llama 4 Scout via Groq with automatic rate-limit retry."""
    if not GROQ_API_KEY:
        raise RuntimeError("GROQ_API_KEY not set in .env")
    import time as _time
    from openai import OpenAI, RateLimitError
    client = OpenAI(api_key=GROQ_API_KEY, base_url="https://api.groq.com/openai/v1")
    for attempt in range(5):
        try:
            raw_resp = client.chat.completions.with_raw_response.create(
                model="meta-llama/llama-4-scout-17b-16e-instruct",
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text",      "text": prompt},
                        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{image_b64}"}},
                    ],
                }],
                max_tokens=max_tokens,
                temperature=0.1,
            )
            response = raw_resp.parse()
            _track_vision_groq_usage(response.usage, raw_resp.headers)
            return response.choices[0].message.content
        except RateLimitError:
            wait = 20 * (attempt + 1)   # 20s, 40s, 60s, 80s, 100s
            print(f"  [vision] Groq rate limit — waiting {wait}s (attempt {attempt+1}/5)")
            _time.sleep(wait)
    raise RuntimeError("Groq vision rate limit exceeded after 5 retries")


def analyze_image_with_llava(image_b64, context_hint="", doc_mode=False, uml_mode=False, fast=False, provider: str | None = None):
    """Send image to LLaVA and extract ALL content useful for Q&A.

    doc_mode=True  → optimised for scanned documents / invoices / tables:
                     verbatim transcription, every table row, all numbers.
    uml_mode=True  → optimised for UML / architecture / ER diagrams:
                     entities, relationships, methods, diagram type.
    doc_mode=False → optimised for photos / screenshots:
                     visual description, colours, structure, purpose.
    """
    if doc_mode:
        prompt = f"""You are a document transcription assistant. Your job is to extract every piece of information from this scanned document page so it can be searched and queried later.

{"Context: " + context_hint if context_hint else ""}

Follow these rules strictly:

1. TRANSCRIBE ALL TEXT verbatim, exactly as it appears — names, dates, addresses, reference numbers, totals. Do not paraphrase.

2. FOR EVERY TABLE, reproduce it completely row by row:
   - Write each row on its own line.
   - Separate columns with " | ".
   - Include the header row first.
   - Count and state the total number of data rows at the end.
   Example:
   CODE | DESCRIPTION | QTY | UNIT PRICE | TOTAL
   A001 | Widget X    | 10  | 5.00       | 50.00
   A002 | Widget Y    | 3   | 12.00      | 36.00
   Total rows: 2

3. LIST ALL NUMBERS AND AMOUNTS exactly as shown: quantities, unit prices, subtotals, taxes, grand totals.

4. STATE KEY FIELDS explicitly:
   - Supplier / vendor name
   - Client / buyer name
   - Document date and reference number
   - Currency and payment terms (if visible)

5. If text is partially illegible, write [illegible] in that spot — do not guess.

Be exhaustive. This transcription is the ONLY source of information for answering user questions about this document."""
    elif uml_mode:
        prompt = f"""Extract all technical information from this diagram. Use compact format to fit everything.

{"File: " + context_hint if context_hint else ""}

Output in this exact format:

DIAGRAM_TYPE: <type>
PURPOSE: <one sentence>

ENTITIES:
<name>: <attr1>, <attr2>, <attr3>
<name>: <attr1>, <attr2>
(one line per entity, list ALL entities visible, do not stop early)

RELATIONSHIPS:
<source> <cardinality> -> <cardinality> <target>: <label>
(one line per relationship, list ALL arrows/connections)

Rules:
- List EVERY entity/class/table visible, scan the entire diagram
- List EVERY arrow or connection
- Use exact names as written in the diagram
- Do not skip any entity even if it seems minor"""
    else:
        prompt = f"""You are a document analysis assistant. Analyze this image completely and extract ALL useful information.

{"Context: " + context_hint if context_hint else ""}

Follow these steps in order:

STEP 1 — TEXT: Copy ALL visible text exactly as written, word for word.

STEP 2 — COLORS: For every distinct element, state its exact color (use specific names: teal, coral, amber, etc.).

STEP 3 — STRUCTURE: Describe the layout and what each element represents.

STEP 4 — TYPE-SPECIFIC details:
- UML diagram: list all classes/actors, attributes, methods, relationships
- Architecture diagram: every component, connection, data flow
- Table: reproduce all rows and columns
- Chart/graph: describe axes, values, legend, each data series and its color
- Screenshot: all visible UI elements, buttons, text, data
- Photo: describe people, setting, actions, notable details
- Logo/illustration: describe every visual element and what it depicts

STEP 5 — PURPOSE: What is the overall meaning or purpose of this image?

Be exhaustive — your description will be the ONLY source of information about this image when answering user questions."""

    if uml_mode:
        num_predict = 3500
    elif doc_mode:
        num_predict = 3000
    else:
        num_predict = 1500

    effective_provider = provider if provider is not None else _active_provider
    if effective_provider == "cloud":
        print(f"  [vision] using Groq Llama 4 Scout")
        return _analyze_image_cloud(image_b64, prompt, max_tokens=num_predict)

    response = ollama.chat(
        model=VISION_MODEL,
        messages=[{
            "role": "user",
            "content": prompt,
            "images": [image_b64]
        }],
        options={
            "num_predict": num_predict,
            "temperature": 0.1,
            "num_gpu": 99,
        }
    )
    return response["message"]["content"]


# ---------------------------------------------------------------------------
# CID ligature map: covers the most common unresolved glyphs pdfplumber
# produces when a PDF embeds a custom font with ligatures.
# ---------------------------------------------------------------------------
_CID_MAP: dict[str, str] = {
    "(cid:11)": "fi", "(cid:12)": "fl", "(cid:13)": "ff",
    "(cid:14)": "ffi", "(cid:15)": "ffl",
    "(cid:28)": "fi", "(cid:29)": "fl", "(cid:30)": "ff",
    "(cid:31)": "ffi", "(cid:32)": "ffl",
    # other common cid glyphs in LaTeX/TeX-generated PDFs
    "(cid:1)":  "!", "(cid:2)": "\"", "(cid:3)": "#",
}

# Characters that appear when pdfplumber mis-maps accented Latin-1 chars
# through a MacRoman or custom encoding.  Key = wrong glyph, value = correct.
_ENCODING_FIX: dict[str, str] = {
    "Ø": "é",   # 0xC3 / MacRoman mismatch — most common in French PDFs
    "ø": "è",
    "Æ": "à",
    "æ": "â",
    "Å": "ê",
    "å": "ë",
    "Ã": "î",
    "ã": "ï",
    "Œ": "ô",
    "œ": "ù",
    "Ç": "ç",
    "ß": "û",
    "†": "°",
}

def _fix_pdf_text(text: str) -> str:
    """
    Post-process pdfplumber output to fix two common artefacts:
    1. (cid:N) sequences — unresolved ligatures / special glyphs.
    2. Wrong characters from MacRoman / custom PDF font encoding.
    """
    for cid, replacement in _CID_MAP.items():
        text = text.replace(cid, replacement)

    text = re.sub(r'\(cid:\d+\)', '', text)

    for wrong, correct in _ENCODING_FIX.items():
        text = text.replace(wrong, correct)

    return text


def extract_pdf_content(file_path, filename, on_progress=None, provider: str = "local"):
    full_content = f"Document: {filename}\n\n"

    with pdfplumber.open(file_path) as pdf:
        total_pages = len(pdf.pages)
        print(f"Processing PDF: {filename} ({total_pages} pages)")
        if on_progress:
            on_progress(0, total_pages)

        for i, page in enumerate(pdf.pages):
            if filename in cancelled_files:
                print(f"Indexing cancelled: {filename}")
                raise InterruptedError(f"Cancelled by user")

            page_num = i + 1
            full_content += f"\n{'='*40}\nPage {page_num}/{total_pages}\n{'='*40}\n"

            text = _fix_pdf_text(page.extract_text() or "")
            if text.strip():
                full_content += f"\n[Text content]\n{text}\n"

            text_len = len(text.strip())

            # Only call LLaVA when pdfplumber extracted nothing at all —
            # i.e. a truly scanned / image-only page with no text layer.
            # LLaVA is NOT reliable as an OCR fallback and hallucinates on
            # text-heavy pages even when images are present (logos, dividers, etc.)
            needs_llava = text_len < 50

            if needs_llava:
                if filename in cancelled_files:
                    raise InterruptedError(f"Cancelled by user")
                try:
                    # Local LLaVA has a 4096-token context window — low resolution keeps image tokens under the limit.
                    # Cloud vision has no such restriction.
                    resolution = 200 if provider == "cloud" else 96
                    page_image = page.to_image(resolution=resolution).original
                    image_b64 = pil_image_to_base64(page_image)
                    visual = analyze_image_with_llava(
                        image_b64,
                        context_hint=f"Page {page_num}/{total_pages} of PDF '{filename}'",
                        doc_mode=True,
                        provider=provider,
                    )
                    full_content += f"\n[Visual content — page {page_num}]\n{visual}\n"
                except InterruptedError:
                    raise
                except Exception as e:
                    print(f"  Page {page_num}: LLaVA error — {e}")
            else:
                print(f"  Page {page_num}: sufficient text ({text_len} chars) — skipping LLaVA")

            if on_progress:
                on_progress(page_num, total_pages)

    return full_content


def extract_image_content(file_path, filename, provider: str = "local"):
    """Extract content from standalone image using LLaVA"""
    print(f"Analyzing image with LLaVA: {filename}")
    image_b64 = image_to_base64(file_path)
    is_uml = _is_uml_image(filename)
    if is_uml:
        print(f"  UML/diagram detected — using diagram extraction prompt")
    description = analyze_image_with_llava(
        image_b64,
        context_hint=f"Image file '{filename}'",
        uml_mode=is_uml,
        provider=provider,
    )
    return f"Image file: {filename}\n\n{description}"


def extract_txt_content(file_path, filename):
    """Read plain text file directly."""
    print(f"Reading text file: {filename}")
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    return f"Text file: {filename}\n\n{text}"


def extract_pptx_content(file_path, filename, provider: str = "local"):
    """Extract text and images from PowerPoint slides — each slide becomes its own labelled block."""
    print(f"Reading PPTX file: {filename}")
    prs = PptxPresentation(file_path)
    total = len(prs.slides)

    def _slide_title(slide):
        """Return title placeholder text, or first non-empty text as fallback."""
        first_text = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            try:
                if shape.is_placeholder and shape.placeholder_format.idx in (0, 1):
                    return text   # proper title placeholder
            except Exception:
                pass
            if not first_text:
                first_text = text  # keep first text as fallback
        return first_text

    slide_titles = []
    slide_first_body = []
    for slide in prs.slides:
        title = _slide_title(slide)
        slide_titles.append(title)
        first_body = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            t = shape.text_frame.text.strip()
            if not t or t == title:
                continue
            first_line = next((ln.strip() for ln in t.splitlines() if ln.strip()), "")
            if first_line:
                first_body = first_line
                break
        slide_first_body.append(first_body)

    # Ordinal words for the first 20 slides so "first slide" matches "Slide 1"
    _ORDINALS = ["first","second","third","fourth","fifth","sixth","seventh",
                 "eighth","ninth","tenth","eleventh","twelfth","thirteenth",
                 "fourteenth","fifteenth","sixteenth","seventeenth","eighteenth",
                 "nineteenth","twentieth"]

    pname = os.path.splitext(filename)[0]

    def _shape_text(shape) -> str:
        """Extract all text from a shape, including tables."""
        if shape.has_text_frame:
            return shape.text_frame.text.strip()
        try:
            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    cells = [c.text_frame.text.strip() for c in row.cells if c.text_frame.text.strip()]
                    if cells:
                        rows.append(" | ".join(cells))
                return "\n".join(rows)
        except Exception:
            pass
        return ""

    # Extract slide 1 body text — title slides carry project name, team, institution
    slide1_body_lines = []
    if prs.slides:
        first_slide = prs.slides[0]
        first_title = slide_titles[0]
        for shape in first_slide.shapes:
            t = _shape_text(shape)
            if t and t != first_title:
                slide1_body_lines.append(t)

    overview_lines = [
        f"  Slide {i} ({_ORDINALS[i-1] if i <= len(_ORDINALS) else ''} slide): {title}" if title
        else f"  Slide {i} ({_ORDINALS[i-1] if i <= len(_ORDINALS) else ''} slide)"
        for i, title in enumerate(slide_titles, 1)
    ]
    slide1_body_text = "\n".join(slide1_body_lines)
    overview = (
        f"Presentation name / title: {pname}\n"
        f"Presentation file: {filename}\n"
        f"Total slides: {total}\n"
        + (f"Cover slide content (project name, team, institution):\n{slide1_body_text}\n" if slide1_body_text else "")
        + f"Slide index (what each slide represents):\n" + "\n".join(overview_lines) + "\n"
    )
    parts = [overview]

    for i, slide in enumerate(prs.slides, 1):
        title_text = slide_titles[i - 1]
        body_lines = []

        for shape in slide.shapes:
            text = _shape_text(shape)
            if not text or text == title_text:
                continue
            body_lines.append(text)

        notes_text = ""
        try:
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass

        # Extract images from slide shapes and analyze with vision
        image_descriptions = []
        try:
            from pptx.util import Pt
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img_bytes = shape.image.blob
                        img_b64 = base64.b64encode(img_bytes).decode("utf-8")
                        desc = analyze_image_with_llava(
                            img_b64,
                            context_hint=f"Image on slide {i} of '{filename}'",
                            provider=provider,
                        )
                        if desc:
                            image_descriptions.append(desc)
                    except Exception as e:
                        print(f"  [pptx] slide {i} image error: {e}")
        except Exception as e:
            print(f"  [pptx] slide {i} image extraction error: {e}")

        ordinal = _ORDINALS[i-1] if i <= len(_ORDINALS) else f"{i}th"
        block = f"\n--- Slide {i} ({ordinal} slide)"
        if title_text:
            block += f": {title_text}"
        block += " ---\n"
        if title_text:
            block += f"Title: {title_text}\n"
        if body_lines:
            block += "\n".join(body_lines) + "\n"
        if image_descriptions:
            block += "\n".join(f"[Image]: {d}" for d in image_descriptions) + "\n"
        if notes_text:
            block += f"Speaker notes: {notes_text}\n"

        parts.append(block)

    return "\n".join(parts)


def extract_uml_content(file_path, filename):
    """Parse PlantUML/UML file into RAG-friendly structured text with English annotations."""
    print(f"Parsing UML file: {filename}")
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        raw = f.read()

    lines = raw.splitlines()

    class_re  = re.compile(r'^\s*class\s+(\w+)', re.IGNORECASE)
    entity_re = re.compile(r'^\s*entity\s+(\w+)', re.IGNORECASE)
    classes = []
    for line in lines:
        m = class_re.match(line) or entity_re.match(line)
        if m and m.group(1) not in classes:
            classes.append(m.group(1))

    class_blocks: dict = {}
    current_class = None
    brace_depth = 0
    for line in lines:
        m = class_re.match(line) or entity_re.match(line)
        if m:
            current_class = m.group(1)
            class_blocks[current_class] = []
            brace_depth = line.count('{') - line.count('}')
            continue
        if current_class:
            brace_depth += line.count('{') - line.count('}')
            if brace_depth <= 0:
                current_class = None
                brace_depth = 0
            else:
                stripped = line.strip().lstrip('+-#~')
                if stripped and stripped not in ('{', '}'):
                    class_blocks[current_class].append(stripped)

    rel_re = re.compile(
        r'^\s*(\w+)\s+"?([^"]*)"?\s+(-+>|<-+|\.+>|<\.+|--)\s+"?([^"]*)"?\s+(\w+)'
        r'(?:\s*:\s*(.+))?'
    )
    rel_simple = re.compile(r'^\s*(\w+)\s+(?:--|-->|\.\.|\.\.>)\s+(\w+)(?:\s*:\s*(.+))?')
    relationships = []
    for line in lines:
        if line.strip().startswith("'"):
            continue
        m = rel_re.match(line)
        if m:
            src, card1, _, card2, dst, label = m.groups()
            label = label.strip() if label else ""
            card1 = (card1 or "").strip()
            card2 = (card2 or "").strip()
            rel_text = f"{src} {label} {dst}" if label else f"{src} is related to {dst}"
            rel_text += f" (cardinality: {card1} to {card2})" if card1 or card2 else ""
            relationships.append(rel_text)
            continue
        m2 = rel_simple.match(line)
        if m2:
            src, dst, label = m2.groups()
            label = (label or "").strip() or "is related to"
            relationships.append(f"{src} {label} {dst}")

    entity_rels: dict = {cls: [] for cls in classes}
    for rel in relationships:
        for cls in classes:
            if re.search(r'\b' + re.escape(cls) + r'\b', rel, re.IGNORECASE):
                entity_rels[cls].append(rel)

    parts = []

    parts.append(
        f"UML Diagram: {filename}\n"
        f"Total entities/tables: {len(classes)}\n"
        f"List of all entities: {', '.join(classes)}\n"
        f"Note: entity names and relationship labels may be in French.\n"
        f"Key translations: Employe=Employee, Salle=Room, Membre=Member, "
        f"Cours=Course, Session=Session, Vestiaire=Locker room, Casier=Locker, "
        f"Equipement=Equipment, Passage=Entry/Access, Souscription=Subscription, "
        f"Reservation=Booking, Programme=Program, Produit=Product, "
        f"Evenement=Event, Maintenance=Maintenance, Role=Role.\n"
        f"Relationship labels: affecte a=assigned to, situe dans=located in, "
        f"se deroule dans=takes place in, contient=contains, loue=rents, "
        f"souscrit=subscribes, effectue=performs, supervise=supervises, "
        f"possede=has, anime=runs, encadre=coaches.\n"
    )

    for cls in classes:
        block = f"\nEntity: {cls}\n"
        attrs = class_blocks.get(cls, [])
        if attrs:
            block += "Attributes:\n" + "\n".join(f"  - {a}" for a in attrs) + "\n"
        rels = entity_rels.get(cls, [])
        if rels:
            block += f"Relationships involving {cls}:\n"
            block += "\n".join(f"  - {r}" for r in rels) + "\n"
        parts.append(block)

    parts.append(f"\n--- Raw UML source ---\n{raw}")

    return "\n".join(parts)


def extract_docx_content(file_path, filename):
    """Extract text and tables from a Word document."""
    print(f"Reading Word document: {filename}")
    doc = DocxDocument(file_path)

    table_summaries = []
    table_blocks = []
    for i, table in enumerate(doc.tables):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            line = " | ".join(c for c in cells if c)
            if line:
                rows.append(line)
        if not rows:
            continue
        first_cell = table.rows[0].cells[0].text.strip() if table.rows else ""
        has_header = first_cell and not first_cell.replace(",", "").replace(".", "").isdigit()
        data_rows = len(rows) - 1 if has_header else len(rows)
        table_summaries.append(f"  Table {i + 1}: {data_rows} data row(s)")
        table_blocks.append((i + 1, data_rows, rows))

    # Only count tables with 2+ data rows as "real" data tables; single-row
    # tables are usually layout/header artefacts in Word documents.
    real_tables = [(t_num, dr, rows) for t_num, dr, rows in table_blocks if dr >= 2]
    total_items = sum(dr for _, dr, _ in real_tables)

    header = f"Document: {filename}\n"
    if table_blocks:
        if real_tables:
            header += f"This document contains {len(real_tables)} data table(s) with {total_items} item(s) in total.\n"
        header += "Table breakdown:\n" + "\n".join(table_summaries) + "\n"

    parts = [header]

    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)

    for t_num, data_rows, rows in table_blocks:
        parts.append(f"\n[Table {t_num}] — {data_rows} data row(s)")
        parts.extend(rows)

    return "\n".join(parts)


def extract_excel_content(file_path, filename):
    """Extract content from an Excel workbook (.xlsx / .xls) as readable text."""
    print(f"Reading Excel file: {filename}")

    def _row_to_str(cells: list) -> str:
        """Convert a row of values to a readable pipe-separated string."""
        str_cells = [str(c).strip() if c is not None else "" for c in cells]
        if not any(str_cells):
            return ""
        # If first column is empty but later columns have values, label as "Total"
        if not str_cells[0] and any(str_cells[1:]):
            str_cells[0] = "Total"
        return " | ".join(str_cells)

    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        use_openpyxl = True
    except Exception:
        wb = None
        use_openpyxl = False

    if not use_openpyxl:
        try:
            import xlrd
            _xlrd_wb = xlrd.open_workbook(file_path)
        except Exception as e:
            return f"Document: {filename}\n[Could not read Excel file: {e}]"

    parts = [f"Document: {filename}\n"]

    if use_openpyxl:
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            if ws.max_row == 0:
                continue
            sheet_rows = []
            for row in ws.iter_rows(values_only=True):
                line = _row_to_str(list(row))
                if line:
                    sheet_rows.append(line)
            if not sheet_rows:
                continue
            data_rows = len(sheet_rows) - 1 if len(sheet_rows) > 1 else len(sheet_rows)
            parts.append(f"\n[Sheet: {sheet_name}] — {data_rows} data row(s)")
            parts.extend(sheet_rows)
    else:
        for i in range(_xlrd_wb.nsheets):
            ws = _xlrd_wb.sheet_by_index(i)
            sheet_rows = []
            for rx in range(ws.nrows):
                cells = [ws.cell_value(rx, cx) for cx in range(ws.ncols)]
                line = _row_to_str(cells)
                if line:
                    sheet_rows.append(line)
            if not sheet_rows:
                continue
            data_rows = len(sheet_rows) - 1 if len(sheet_rows) > 1 else len(sheet_rows)
            parts.append(f"\n[Sheet: {ws.name}] — {data_rows} data row(s)")
            parts.extend(sheet_rows)

    return "\n".join(parts)


def get_nodes_for_files(file_names: list) -> list:
    """Fetch text nodes from ChromaDB so BM25 can index them for keyword search."""
    nodes = []
    for fname in file_names:
        try:
            results = chroma_collection.get(
                where={"file_name": fname},
                include=["documents", "metadatas"]
            )
            for doc_id, text, meta in zip(
                results["ids"], results["documents"], results["metadatas"]
            ):
                nodes.append(TextNode(id_=doc_id, text=text or "", metadata=meta or {}))
        except Exception as e:
            print(f"BM25 node fetch error for {fname}: {e}")
    return nodes



def _looks_like_mld(text: str) -> bool:
    """Heuristic: text contains ≥3 relational-schema entity definitions like: word = (...)."""
    return len(re.findall(r'\w+\s*=\s*\(', text)) >= 3


def _split_schema_by_semicolons(text: str, filename: str, doc_type: str) -> list:
    """
    For relational-schema (MLD) text where each entity definition ends with ';',
    produce one TextNode per entity instead of splitting by token count.
    Also prepends an overview node listing all entity/table names so that
    exhaustive queries ("list all entities") always retrieve the full index
    regardless of top-K.
    """
    entity_nodes = []
    entity_names = []

    for part in text.split(';'):
        clean = part.strip()
        # Drop empty parts and short extractor-added headers (page markers, etc.)
        if not clean or len(clean) < 15:
            continue
        # Strip page-marker lines (===...===) that may be prepended by the extractor
        lines = [l for l in clean.splitlines()
                 if not l.strip().startswith('=') and l.strip() not in ('[Text content]', '')]
        clean = '\n'.join(lines).strip()
        if not clean or len(clean) < 10:
            continue
        # Grab the entity/table name: last word before the first '='
        entity_name = ''
        if '=' in clean:
            candidate = clean.split('=')[0].strip().split()
            entity_name = candidate[-1] if candidate else ''
            if entity_name:
                entity_names.append(entity_name)
        entity_nodes.append(TextNode(
            text=clean + ';',
            metadata={
                'file_name': filename,
                'doc_type':  doc_type,
                'entity':    entity_name,
                'schema_chunk': 'entity',
            }
        ))

    # Overview node — lists every table/entity name in one chunk.
    # Scores highest on "list all entities / tables" queries and ensures
    # exhaustive answers even when top-K < total entity count.
    if entity_names:
        overview_text = (
            f"MLD schema overview for {filename}\n"
            f"Total tables/entities: {len(entity_names)}\n"
            f"Entity list: {', '.join(entity_names)}\n"
            f"Tables: {'; '.join(entity_names)}"
        )
        overview_node = TextNode(
            text=overview_text,
            metadata={
                'file_name': filename,
                'doc_type':  doc_type,
                'entity':    'overview',
                'schema_chunk': 'overview',
            }
        )
        return [overview_node] + entity_nodes

    return entity_nodes

def add_document_to_index(file_path, filename, provider: str | None = None):
    global index
    if provider is None:
        provider = _active_provider  # snapshot at submission time
    try:
        indexing_status[filename] = "indexing"
        indexing_progress[filename] = {"current": 0, "total": 0}
        extension = filename.lower().split('.')[-1]

        if extension == 'pdf':
            def _on_progress(current, total):
                indexing_progress[filename] = {"current": current, "total": total}
            text = extract_pdf_content(file_path, filename, on_progress=_on_progress, provider=provider)
            doc_type = "pdf"
        elif extension in ['png', 'jpg', 'jpeg', 'gif', 'bmp', 'webp']:
            if filename in cancelled_files:
                raise InterruptedError("Cancelled by user")
            text = extract_image_content(file_path, filename, provider=provider)
            doc_type = "image"
        elif extension == 'txt':
            text = extract_txt_content(file_path, filename)
            doc_type = "txt"
        elif extension == 'docx':
            text = extract_docx_content(file_path, filename)
            doc_type = "docx"
        elif extension in ['xlsx', 'xls']:
            text = extract_excel_content(file_path, filename)
            doc_type = "excel"
        elif extension == 'pptx':
            text = extract_pptx_content(file_path, filename, provider=provider)
            doc_type = "pptx"
        elif extension in ['puml', 'plantuml', 'uml']:
            text = extract_uml_content(file_path, filename)
            doc_type = "uml"
        elif extension in ['md', 'csv']:
            text = extract_txt_content(file_path, filename)
            doc_type = "txt"
        else:
            text = f"[Unsupported file: {filename}]"
            doc_type = "unknown"

        # Relational schemas (MLD / SQL DDL) have a natural delimiter: ';'.
        # Fixed-size token splitting cuts entities in half, making it impossible
        # for the LLM to read a definition in one piece.
        # When the extracted text looks like an MLD, bypass HierarchicalNodeParser
        # and create one TextNode per entity instead.
        if _looks_like_mld(text):
            print(f"Schema detected in {filename} — using semicolon-delimiter chunking")
            schema_nodes = _split_schema_by_semicolons(text, filename, doc_type)
            print(f"  → {len(schema_nodes)} entity chunks")

            # Clear stale nodes from docstore AND ChromaDB before inserting new ones
            try:
                stale = chroma_collection.get(where={"file_name": filename})
                if stale["ids"]:
                    chroma_collection.delete(ids=stale["ids"])
                    print(f"  Deleted {len(stale['ids'])} stale ChromaDB chunks for {filename}")
            except Exception as e:
                print(f"  Warning: could not clear stale ChromaDB chunks: {e}")

            if PARENT_CHILD_AVAILABLE:
                old_node_ids = [
                    nid for nid, node in storage_context.docstore.docs.items()
                    if node.metadata.get("file_name") == filename
                ]
                for nid in old_node_ids:
                    try:
                        storage_context.docstore.delete_document(nid)
                    except Exception:
                        pass

            if index is None:
                index = VectorStoreIndex(
                    schema_nodes, storage_context=storage_context, show_progress=False
                )
            else:
                index.insert_nodes(schema_nodes)

        else:
            doc = Document(
                text=text,
                metadata={"file_name": filename, "file_path": file_path, "doc_type": doc_type}
            )

            if PARENT_CHILD_AVAILABLE:
                # Two levels: parent (~512 tokens, in docstore) + leaf (~128 tokens,
                # in ChromaDB). AutoMergingRetriever promotes siblings to parent at
                # query time → LLM gets richer context, retrieval stays precise.

                # Remove stale docstore nodes for this file before re-indexing so
                # ChromaDB and the docstore never reference each other's old IDs.
                old_node_ids = [
                    nid for nid, node in storage_context.docstore.docs.items()
                    if node.metadata.get("file_name") == filename
                ]
                for nid in old_node_ids:
                    try:
                        storage_context.docstore.delete_document(nid)
                    except Exception:
                        pass
                if old_node_ids:
                    print(f"Removed {len(old_node_ids)} stale docstore nodes for {filename}")

                parser     = HierarchicalNodeParser.from_defaults(
                    chunk_sizes=[PARENT_CHUNK_SIZE, CHILD_CHUNK_SIZE]
                )
                all_nodes  = parser.get_nodes_from_documents([doc])
                leaf_nodes = get_leaf_nodes(all_nodes)

                for node in all_nodes:
                    node.metadata.setdefault("file_name", filename)
                    node.metadata.setdefault("doc_type",  doc_type)

                storage_context.docstore.add_documents(all_nodes)
                _persist_docstore()

                if index is None:
                    index = VectorStoreIndex(
                        leaf_nodes, storage_context=storage_context, show_progress=False
                    )
                else:
                    index.insert_nodes(leaf_nodes)
            else:
                if index is None:
                    index = VectorStoreIndex.from_documents([doc], storage_context=storage_context)
                else:
                    index.insert(doc)

        indexing_status[filename] = "ready"
        indexing_progress.pop(filename, None)
        cancelled_files.discard(filename)
        print(f"Ready: {filename} ({doc_type}, {len(text)} chars)")

    except InterruptedError:
        indexing_status[filename] = "cancelled"
        indexing_progress.pop(filename, None)
        cancelled_files.discard(filename)
        # delete the physical file since indexing was cancelled
        if os.path.exists(file_path):
            os.remove(file_path)
        print(f"Cancelled and cleaned: {filename}")

    except Exception as e:
        indexing_status[filename] = "error"
        indexing_progress.pop(filename, None)
        print(f"Indexing error for {filename}: {e}")

_HEADERS = {"User-Agent": "Mozilla/5.0 (compatible; RAGBot/1.0)"}
_SKIP_TAGS = {"script", "style", "nav", "footer", "header", "aside",
              "noscript", "form", "button", "iframe", "svg"}

def _fetch_url_text(url: str) -> tuple[str, str]:
    """Return (title, clean_text) extracted from a URL."""
    resp = requests.get(url, headers=_HEADERS, timeout=15)
    resp.raise_for_status()
    soup = BeautifulSoup(resp.text, "html.parser")

    for tag in soup(_SKIP_TAGS):
        tag.decompose()

    title = soup.title.string.strip() if soup.title and soup.title.string else url

    body = (soup.find("article") or soup.find("main") or soup.find("body"))
    lines = []
    if body:
        for el in body.find_all(["h1","h2","h3","h4","p","li","td","th","blockquote"]):
            t = el.get_text(" ", strip=True)
            if t:
                lines.append(t)
    text = "\n".join(lines) if lines else soup.get_text("\n", strip=True)

    return title, text

def _safe_filename(title: str, url: str) -> str:
    """Turn a page title into a safe .txt filename, fallback to domain."""
    name = re.sub(r'[\\/*?:"<>|]', "", title)[:80].strip()
    if not name:
        from urllib.parse import urlparse
        name = urlparse(url).netloc
    return name + ".txt"


@app.post("/upload-url")
async def upload_url(req: UrlIngestRequest,
                     current_user: UserModel = Depends(get_current_user)):
    """Fetch a web URL, extract clean text, and index it like any uploaded file."""
    try:
        title, text = _fetch_url_text(req.url)
    except requests.exceptions.RequestException as e:
        raise HTTPException(400, f"Could not fetch URL: {e}")

    if len(text.strip()) < 100:
        raise HTTPException(400, "Page returned too little text — it may require JavaScript to render.")

    filename = _safe_filename(title, req.url)
    dest = os.path.join(UPLOAD_DIR, filename)

    with open(dest, "w", encoding="utf-8") as f:
        f.write(f"Source: {req.url}\n\n{text}")

    file_owners[filename] = current_user.id
    _save_file_owners()

    indexing_status[filename]   = "indexing"
    indexing_progress[filename] = {"current": 0, "total": 0}
    cancelled_files.discard(filename)
    loop = asyncio.get_event_loop()
    loop.run_in_executor(executor, add_document_to_index, dest, filename, _active_provider)

    return {"name": filename, "status": "indexing", "title": title}


@app.post("/upload")
async def upload(file: UploadFile = File(...),
                 current_user: UserModel = Depends(get_current_user)):
    """
    Save file and start indexing in background.
    If the same file (identical content) is re-uploaded, returns "ready" immediately.
    Returns immediately — frontend polls /status/{filename} to know when ready.
    """
    content = await file.read()
    incoming_hash = hashlib.md5(content).hexdigest()

    # Strip repeated extensions e.g. "résumé.pdf.pdf" → "résumé.pdf"
    _stem, _ext = os.path.splitext(file.filename)
    if os.path.splitext(_stem)[1].lower() == _ext.lower():
        file.filename = _stem

    # Skip re-indexing if the file is already indexed and content is unchanged
    if (
        file_hashes.get(file.filename) == incoming_hash
        and indexing_status.get(file.filename) == "ready"
    ):
        # Re-claim ownership (user may have refreshed after backend restart)
        file_owners[file.filename] = current_user.id
        _save_file_owners()
        print(f"Skipping re-index (unchanged): {file.filename}")
        return {"id": file.filename, "name": file.filename, "status": "ready"}

    path = f"{UPLOAD_DIR}/{file.filename}"

    try:
        results = chroma_collection.get(where={"file_name": file.filename})
        if results["ids"]:
            chroma_collection.delete(ids=results["ids"])
            print(f"Cleaned old chunks for {file.filename}")
    except Exception as e:
        print(f"Cleanup error: {e}")

    os.makedirs(UPLOAD_DIR, exist_ok=True)
    with open(path, "wb") as f:
        f.write(content)

    file_owners[file.filename] = current_user.id
    _save_file_owners()
    file_hashes[file.filename] = incoming_hash
    indexing_status[file.filename] = "indexing"
    loop = asyncio.get_event_loop()
    loop.run_in_executor(executor, add_document_to_index, path, file.filename, _active_provider)

    # Pre-convert PPTX → PDF in the background so the preview is ready immediately
    if file.filename.lower().endswith(".pptx"):
        loop.run_in_executor(executor, _pptx_to_pdf_cached, path, file.filename)

    return {"id": file.filename, "name": file.filename, "status": "indexing"}

@app.get("/status/{filename}")
def get_status(filename: str,
               current_user: UserModel = Depends(get_current_user)):
    """
    Poll this endpoint to know when a file is ready.
    Returns: indexing | ready | error | unknown
    """
    status   = indexing_status.get(filename, "unknown")
    progress = indexing_progress.get(filename)  # {"current": int, "total": int} or None
    return {"filename": filename, "status": status, "progress": progress}


@app.get("/documents")
def list_documents(current_user: UserModel = Depends(get_current_user)):
    all_files = [f for f in os.listdir(UPLOAD_DIR) if not f.startswith('.')]
    # Admins see everything; regular users see only their own files
    if current_user.role != "admin":
        all_files = [f for f in all_files
                     if file_owners.get(f) == current_user.id]
    return [
        {"id": f, "name": f, "status": indexing_status.get(f, "ready"),
         "owner": file_owners.get(f)}
        for f in all_files
    ]


@app.get("/index/stats")
def index_stats(current_user: UserModel = Depends(_require_admin)):
    return {
        "total_chunks": chroma_collection.count(),
        "files": os.listdir(UPLOAD_DIR),
        "indexing_status": indexing_status
    }


@app.get("/debug/chunks/{filename}")
def debug_chunks(filename: str, current_user: UserModel = Depends(_require_admin)):
    """Return all stored chunk texts for a file — for debugging extraction."""
    results = chroma_collection.get(where={"file_name": filename}, include=["documents"])
    chunks = results.get("documents", [])
    return {"filename": filename, "count": len(chunks), "chunks": chunks}


@app.get("/dashboard")
def dashboard(current_user: UserModel = Depends(get_current_user)):
    files = [f for f in os.listdir(UPLOAD_DIR) if not f.startswith('.')]
    ready = [f for f in files if indexing_status.get(f, "ready") == "ready"]
    indexing = [f for f in files if indexing_status.get(f, "") == "indexing"]

    file_chunks = {}
    for f in files:
        try:
            res = chroma_collection.get(where={"file_name": f})
            file_chunks[f] = len(res["ids"])
        except Exception:
            file_chunks[f] = 0

    return {
        "models": {
            "llm":    _last_groq_model if _active_provider == "cloud" else LLM_MODEL,
            "embed":  EMBED_MODEL,
            "vision": "llama-4-scout-17b (Groq)" if _active_provider == "cloud" else VISION_MODEL,
            "provider": _active_provider,
        },
        "documents": {
            "total": len(files),
            "ready": len(ready),
            "indexing": len(indexing),
            "file_chunks": file_chunks,
        },
        "chunks": {
            "total": chroma_collection.count(),
        },
        "config": {
            "similarity_top_k": SIMILARITY_TOP_K,
            "max_upload_mb": MAX_UPLOAD_MB,
        },
        "tokens": {
            "prompt": token_usage["prompt"],
            "completion": token_usage["completion"],
            "total": token_usage["prompt"] + token_usage["completion"],
            "requests": token_usage["requests"],
        },
        "groq_tokens": {
            "date": groq_token_usage.get("date"),
            "models": {
                model: (lambda d=data, m=model: {
                    "prompt":      d["prompt"],
                    "completion":  d["completion"],
                    # Prefer authoritative values from Groq response headers
                    "total":       d.get("used_actual") if d.get("used_actual") is not None else d["prompt"] + d["completion"],
                    "daily_limit": d.get("limit_actual") or GROQ_TPD_LIMITS.get(m, 100_000),
                    "pct": round(
                        (d.get("used_actual") if d.get("used_actual") is not None else d["prompt"] + d["completion"])
                        / max(d.get("limit_actual") or GROQ_TPD_LIMITS.get(m, 100_000), 1) * 100, 1
                    ),
                    "from_headers": d.get("used_actual") is not None,
                    "tpm_limit":     d.get("tpm_limit"),
                    "tpm_remaining": d.get("tpm_remaining"),
                    "tpm_reset":     d.get("tpm_reset"),
                })()
                for model, data in groq_token_usage.get("models", {}).items()
            },
        },
        "queries": {
            "total": query_stats["total"],
            "avg_response_ms": (
                round(query_stats["total_ms"] / query_stats["total"])
                if query_stats["total"] > 0 else 0
            ),
        },
    }


@app.get("/files/{filename}")
def serve_file(filename: str,
               current_user: UserModel = Depends(get_current_user)):
    if current_user.role != "admin" and file_owners.get(filename) != current_user.id:
        raise HTTPException(403, "Access denied")
    path = _safe_upload_path(filename)
    if not os.path.exists(path):
        raise HTTPException(404, "File not found")
    import mimetypes
    mime, _ = mimetypes.guess_type(filename)
    return FileResponse(path, media_type=mime or "application/octet-stream")


@app.get("/preview/{filename}")
def preview_file(filename: str,
                 current_user: UserModel = Depends(get_current_user)):
    """Return a plain-text preview (first 8000 chars of extracted content)."""
    if current_user.role != "admin" and file_owners.get(filename) != current_user.id:
        raise HTTPException(403, "Access denied")
    path = _safe_upload_path(filename)
    if not os.path.exists(path):
        raise HTTPException(404, "File not found")
    ext = filename.lower().rsplit('.', 1)[-1]
    try:
        if ext == 'docx':
            text = extract_docx_content(path, filename)
        elif ext in ['xlsx', 'xls']:
            text = extract_excel_content(path, filename)
        elif ext == 'pptx':
            text = extract_pptx_content(path, filename)
        elif ext in ['puml', 'plantuml', 'uml', 'txt', 'md', 'csv']:
            text = extract_txt_content(path, filename)
        else:
            text = ""
    except Exception as e:
        text = f"[Preview error: {e}]"
    return {"text": text[:8000]}


_LO_BIN: str | None = None
_LO_BIN_CHECKED: bool = False


def _lo_startupinfo():
    """Return a STARTUPINFO that hides the window on Windows; None on other platforms."""
    if os.name == "nt":
        si = subprocess.STARTUPINFO()
        si.dwFlags |= subprocess.STARTF_USESHOWWINDOW
        si.wShowWindow = 0  # SW_HIDE
        return si
    return None


def _find_libreoffice() -> str | None:
    """Return path to a working LibreOffice binary (cached; found by file existence, no subprocess probe)."""
    import glob, shutil
    global _LO_BIN, _LO_BIN_CHECKED
    if _LO_BIN_CHECKED:
        return _LO_BIN

    # 1. Check PATH first (Linux / macOS / Windows with soffice on PATH)
    for name in ("soffice", "libreoffice"):
        found = shutil.which(name)
        if found:
            _LO_BIN = found
            _LO_BIN_CHECKED = True
            return _LO_BIN

    # 2. Glob every versioned LibreOffice install under Program Files (Windows)
    for prog_dir in [r"C:\Program Files", r"C:\Program Files (x86)"]:
        for lo_dir in sorted(glob.glob(os.path.join(prog_dir, "LibreOffice*")), reverse=True):
            exe = os.path.join(lo_dir, "program", "soffice.exe")
            if os.path.isfile(exe):
                _LO_BIN = exe
                _LO_BIN_CHECKED = True
                return _LO_BIN

    _LO_BIN_CHECKED = True  # not found, cache the miss
    return None


def _doc_to_pdf_cached(src: str, filename: str) -> str | None:
    """Convert any LibreOffice-compatible file (PPTX, DOCX, XLSX…) to PDF and cache it.
    Returns cached PDF path, or None on failure."""
    import tempfile
    ext        = os.path.splitext(filename)[1].lower()   # e.g. ".docx"
    stem       = os.path.splitext(filename)[0]
    mtime      = int(os.path.getmtime(src))
    cache_name = f"{stem}_{mtime}.pdf"
    cached     = os.path.join(SLIDES_CACHE_DIR, cache_name)

    if os.path.exists(cached):
        return cached

    lo_bin = _find_libreoffice()
    if not lo_bin:
        return None

    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            safe_src = os.path.join(tmpdir, f"input{ext}")
            shutil.copy2(src, safe_src)
            result = subprocess.run(
                [lo_bin, "--headless", "--norestore", "--nofirststartwizard",
                 "--convert-to", "pdf", "--outdir", tmpdir, safe_src],
                capture_output=True, stdin=subprocess.DEVNULL,
                startupinfo=_lo_startupinfo(), timeout=120,
            )
            if result.returncode != 0:
                return None
            converted = os.path.join(tmpdir, f"input.pdf")
            if not os.path.exists(converted):
                return None
            shutil.move(converted, cached)
    except Exception:
        return None

    return cached


def _pptx_to_pdf_cached(src: str, filename: str) -> str | None:
    """Alias kept for backwards compatibility — delegates to _doc_to_pdf_cached."""
    return _doc_to_pdf_cached(src, filename)


@app.get("/slides-pdf/{filename}")
def slides_pdf(filename: str,
               current_user: UserModel = Depends(get_current_user)):
    """Return a cached PDF rendition of a PPTX file (converts on demand)."""
    if current_user.role != "admin" and file_owners.get(filename) != current_user.id:
        raise HTTPException(403, "Access denied")

    src = _safe_upload_path(filename)
    if not os.path.exists(src):
        raise HTTPException(404, "File not found")

    # _pptx_to_pdf_cached calls _find_libreoffice internally; no need to call it twice
    cached = _pptx_to_pdf_cached(src, filename)
    if not cached:
        if _find_libreoffice() is None:
            raise HTTPException(500,
                "LibreOffice not found. Install it from https://www.libreoffice.org "
                "and make sure 'soffice' is on your PATH.")
        raise HTTPException(500, "LibreOffice conversion failed")

    stem = os.path.splitext(filename)[0]
    return FileResponse(cached, media_type="application/pdf", filename=stem + ".pdf")


_DOC_PREVIEW_EXTS = {'.docx', '.doc', '.xlsx', '.xls', '.odt', '.ods', '.csv'}

@app.get("/doc-pdf/{filename}")
def doc_pdf(filename: str,
            current_user: UserModel = Depends(get_current_user)):
    """Return a PDF rendition of a Word / Excel file for in-browser preview."""
    if current_user.role != "admin" and file_owners.get(filename) != current_user.id:
        raise HTTPException(403, "Access denied")

    ext = os.path.splitext(filename)[1].lower()
    if ext not in _DOC_PREVIEW_EXTS:
        raise HTTPException(400, f"Unsupported file type for PDF preview: {ext}")

    src = _safe_upload_path(filename)
    if not os.path.exists(src):
        raise HTTPException(404, "File not found")

    cached = _doc_to_pdf_cached(src, filename)
    if not cached:
        if _find_libreoffice() is None:
            raise HTTPException(500,
                "LibreOffice not found. Install it to enable Word/Excel preview.")
        raise HTTPException(500, "LibreOffice conversion failed")

    stem = os.path.splitext(filename)[0]
    return FileResponse(cached, media_type="application/pdf", filename=stem + ".pdf")


class TitleRequest(BaseModel):
    question: str
    files: list[str] = []

class ProviderRequest(BaseModel):
    provider: str  # "local" or "cloud"

@app.post("/provider")
async def set_provider(req: ProviderRequest, current_user: UserModel = Depends(get_current_user)):
    global _active_provider
    if req.provider not in ("local", "cloud"):
        raise HTTPException(status_code=400, detail="provider must be 'local' or 'cloud'")
    _active_provider = req.provider
    _save_provider(_active_provider)
    print(f"[provider] switched to {_active_provider}")
    return {"provider": _active_provider}


@app.post("/title")
async def generate_title(req: TitleRequest):
    """Generate a specific 2-5 word session title from the question and file names."""
    files_hint = ""
    if req.files:
        names = ", ".join(req.files)
        files_hint = f"Files involved: {names}\n"
    prompt = (
        "Write a 2 to 5 word title that captures the specific action AND the subject. "
        "Include the document name or key subject so the user knows exactly what this chat is about. "
        "Return ONLY the title — no punctuation, no quotes, no explanation.\n\n"
        f"{files_hint}"
        f"Request: {req.question}\n\nTitle:"
    )
    result = await Settings.llm.acomplete(prompt)
    title = str(result).strip().strip('"\'').strip('.')
    return {"title": title}


_COMPARISON_KEYWORDS = {
    # English
    'compare', 'comparison', 'contrast', 'difference', 'differences',
    'versus', ' vs ', ' vs.', 'similar', 'similarity', 'both documents',
    'both files', 'each document', 'each file', 'across documents',
    'across files', 'which document', 'which file', 'between the two',
    'from the pdf', 'from the uml', 'from both', 'in both',
    'the pdf and', 'and the pdf', 'the uml and', 'and the uml',
    'the diagram and', 'and the diagram',
    # French
    'comparer', 'comparaison', 'différence', 'différences',
    'les deux documents', 'les deux fichiers', 'chaque document',
    'à partir du pdf', 'à partir du diagramme', 'à partir des documents',
    'depuis le pdf', 'depuis le diagramme', 'selon le pdf', 'selon le uml',
    'le pdf et', 'et le pdf', 'le diagramme et', 'et le diagramme',
    'du pdf et', 'et du pdf', 'du diagramme et', 'et du diagramme',
    'dans les deux', 'les fichiers', 'ces documents',
}

def is_comparison_query(question: str, num_files: int) -> bool:
    """Return True when the question likely asks for a cross-document comparison."""
    if num_files < 2:
        return False
    q = question.lower()
    return any(kw in q for kw in _COMPARISON_KEYWORDS)


_EXHAUSTIVE_KEYWORDS = (
    # English
    "list all", "list every", "list each", "show all", "show every",
    "enumerate", "extract all", "extract every", "extract each",
    "give me all", "give me every", "give me a list",
    "what are all", "what are the", "all the",
    "summarize", "summarise", "summary", "resume", "résumé",
    "overview", "full list", "complete list",
    # French
    "liste", "lister", "listez", "liste toutes", "liste tous",
    "énumère", "énumérer", "énumérez",
    "extraire", "extrait", "extraire tout", "extraire toutes",
    "résume", "résumer", "résumé", "récapitule", "récapitulatif",
    "toutes les", "tous les", "l'ensemble",
    "quelles sont", "quels sont",
)

def is_exhaustive_query(question: str) -> bool:
    """Return True when the question asks for a complete list, summary, or extraction."""
    q = question.lower()
    return any(kw in q for kw in _EXHAUSTIVE_KEYWORDS)


async def retrieve_per_file(
    standalone: str,
    question_files: list[str],
    chunks_per_file: int,
) -> tuple[list, str, list[str]]:
    """
    Retrieve `chunks_per_file` chunks from each file independently,
    then assemble a labeled context so the LLM knows what belongs to what.
    Returns (all_nodes, labeled_context, sources).
    """
    all_nodes: list = []
    context_parts: list[str] = []

    loop = asyncio.get_event_loop()

    for fname in question_files:
        file_filter = MetadataFilters(
            filters=[MetadataFilter(key="file_name", value=fname, operator=FilterOperator.EQ)],
            condition=FilterCondition.OR,
        )
        vec_ret = index.as_retriever(similarity_top_k=chunks_per_file, filters=file_filter)

        file_bm25_nodes = get_nodes_for_files([fname])
        if file_bm25_nodes:
            bm25_ret = BM25Retriever.from_defaults(
                nodes=file_bm25_nodes, similarity_top_k=chunks_per_file
            )
            ret = QueryFusionRetriever(
                [vec_ret, bm25_ret],
                similarity_top_k=chunks_per_file,
                num_queries=1,
                mode="reciprocal_rerank",
                use_async=False,
            )
        else:
            ret = vec_ret

        file_nodes = await ret.aretrieve(standalone)

        # Per-file reranking — keep scores balanced across files
        if reranker and len(file_nodes) > 1:
            captured = list(file_nodes)
            q_str    = standalone
            k        = chunks_per_file

            def _rerank_file(nodes_in=captured, q=q_str, top=k):
                pairs  = [(q, n.get_content()) for n in nodes_in]
                scores = reranker.predict(pairs)
                ranked = sorted(zip(nodes_in, scores), key=lambda x: x[1], reverse=True)
                return [n for n, _ in ranked[:top]]

            file_nodes = await loop.run_in_executor(None, _rerank_file)

        if file_nodes:
            all_nodes.extend(file_nodes)
            file_text = "\n\n".join(n.get_content() for n in file_nodes)
            context_parts.append(f"=== Document: {fname} ===\n{file_text}")

    labeled_context = "\n\n".join(context_parts)
    sources = list({n.metadata.get("file_name", "unknown") for n in all_nodes})
    return all_nodes, labeled_context, sources


async def condense_question(question: str, history: list[HistoryEntry], llm=None) -> str:
    """Normalise typos/abbreviations and, when there is history, rewrite as a standalone question."""
    llm = llm or Settings.llm
    if not history:
        prompt = (
            "Fix any spelling mistakes, typos, or grammatical errors in the following question "
            "and rephrase it as a clear, well-formed question. "
            "Preserve the original meaning exactly — do not add or remove information.\n"
            "Return ONLY the corrected question — no explanation.\n\n"
            f"Original: {question}\n\n"
            "Corrected:"
        )
    else:
        history_text = "\n".join(
            f"User: {h.question}\nAssistant: {h.answer}" for h in history[-5:]
        )
        prompt = (
            "Given the conversation below and a follow-up question, do two things:\n"
            "1. Fix any spelling mistakes or typos in the follow-up question.\n"
            "2. Rewrite it as a fully standalone question that includes all necessary context from the history.\n"
            "Return ONLY the rewritten question — no explanation.\n\n"
            f"Conversation:\n{history_text}\n\n"
            f"Follow-up: {question}\n\n"
            "Standalone question:"
        )
    result = await llm.acomplete(prompt)
    record_tokens(result)
    condensed = str(result).strip()
    print(f"[condense] '{question}' → '{condensed}'")
    return condensed


async def hyde_expand(question: str, llm=None) -> str:
    """HyDE: generate a short hypothetical document passage for vector similarity search."""
    llm = llm or Settings.llm
    prompt = (
        "Write a concise passage (3-5 sentences) that would directly answer "
        "the following question if it appeared in a document. "
        "Do not mention that this is hypothetical. Just write the passage.\n\n"
        f"Question: {question}\n\nPassage:"
    )
    try:
        result = await llm.acomplete(prompt)
        record_tokens(result)
        hypothesis = str(result).strip()
        print(f"[hyde] hypothesis: {hypothesis[:120]}…")
        return hypothesis
    except Exception as e:
        print(f"[hyde] failed ({e}), falling back to original query")
        return question


async def multi_query_expand(question: str, n: int = 3, llm=None) -> list[str]:
    """Generate n alternative phrasings of the question for multi-query retrieval."""
    llm = llm or Settings.llm
    prompt = (
        f"Generate {n} different phrasings of the following question to improve document retrieval. "
        "Each should approach the same information need from a different angle "
        "(e.g. different vocabulary, more specific, more general). "
        "Output only the questions, one per line, no numbering or extra text.\n\n"
        f"Question: {question}\n\nAlternative phrasings:"
    )
    try:
        result = await llm.acomplete(prompt)
        record_tokens(result)
        lines = [l.strip().lstrip("•-–1234567890.) ") for l in str(result).strip().splitlines() if l.strip()]
        rewrites = [l for l in lines if l and l != question][:n]
        print(f"[multi-query] rewrites: {rewrites}")
        return rewrites if rewrites else [question]
    except Exception as e:
        print(f"[multi-query] failed ({e}), using original query")
        return [question]


@app.post("/ask")
async def ask(q: Question,
              current_user: UserModel = Depends(get_current_user)):
    if not q.files:
        raise HTTPException(400, "No documents in this chat. Upload a file to get started.")

    if current_user.role != "admin":
        forbidden = [f for f in q.files if file_owners.get(f) != current_user.id]
        if forbidden:
            raise HTTPException(403, f"Access denied to: {', '.join(forbidden)}")

    # Wait inside the SSE generator (not here) so the client can cancel at any time.
    session_indexing = [f for f in q.files if indexing_status.get(f) == "indexing"]

    if not index:
        raise HTTPException(400, "No documents indexed yet. Please upload a file first.")

    # Capture for generator closure
    question       = q.question
    history        = q.history
    question_files = q.files
    global _last_groq_model
    if q.provider == "cloud" and q.groq_model:
        _last_groq_model = q.groq_model
    active_llm     = _get_llm(q.provider, groq_model=q.groq_model)
    # For cloud, route cheap auxiliary calls (condense/HyDE/multi-query/eval) to the
    # 8B model (500K TPD) so the 70B quota is reserved for the actual answer.
    aux_llm        = (_get_llm(q.provider, groq_model=GROQ_AUX_MODEL)
                      if q.provider in ("cloud", "groq") else active_llm)

    def build_citations(nodes: list) -> list[dict]:
        """
        Build structured citations from retrieved nodes.
        Scans each chunk's text for 'Page X/Y' markers (embedded by extract_pdf_content)
        and groups page numbers by filename.
        Returns: [{"file": "foo.pdf", "pages": [1, 3]}, ...]
        Images/docx/txt have no page markers → pages list is empty.
        """
        _page_re = re.compile(r'Page\s+(\d+)/\d+', re.IGNORECASE)
        file_pages: dict[str, set] = {}
        for node in nodes:
            fname = node.metadata.get("file_name", "unknown")
            pages_found = {int(p) for p in _page_re.findall(node.get_content())}
            if fname not in file_pages:
                file_pages[fname] = set()
            file_pages[fname].update(pages_found)
        return [
            {"file": fname, "pages": sorted(pgs)}
            for fname, pgs in file_pages.items()
        ]

    async def event_stream():
        try:
            if session_indexing:
                print(f"[/ask] Waiting for {session_indexing} to finish indexing...")
                wait_timeout = 600
                elapsed = 0
                while any(indexing_status.get(f) == "indexing" for f in question_files):
                    yield f"data: {json.dumps({'type': 'indexing_wait', 'files': session_indexing})}\n\n"
                    await asyncio.sleep(2)
                    elapsed += 2
                    if elapsed >= wait_timeout:
                        yield f"data: {json.dumps({'type': 'error', 'message': 'Indexing timed out. Please try again.'})}\n\n"
                        return
                print("[/ask] Indexing done — proceeding to answer.")

            # Early-exit if we already know the daily quota is exhausted —
            # avoids wasting aux-model tokens on condense/HyDE/multi-query
            # before hitting the same wall on the main answer call.
            if _groq_daily_limit_hit and q.provider in ("cloud", "groq"):
                yield f"data: {json.dumps({'type': 'error', 'message': 'Rate limit reached — Groq daily token quota exhausted. Please try again tomorrow or switch to a different model.'})}\n\n"
                return

            standalone = question if (q.fast or not history) else await condense_question(question, history, llm=aux_llm)

            if _groq_daily_limit_hit and q.provider in ("cloud", "groq"):
                yield f"data: {json.dumps({'type': 'error', 'message': 'Rate limit reached — Groq daily token quota exhausted. Please try again tomorrow or switch to a different model.'})}\n\n"
                return

            comparison_mode = is_comparison_query(question, len(question_files))

            exhaustive = is_exhaustive_query(question)
            if comparison_mode:
                # Raise chunks_per_file generously — structured files (UML, MLD, schema)
                # need ALL entity blocks, and regular comparison benefits from more context too.
                chunks_per_file = max(8, SIMILARITY_TOP_K)
                print(f"[compare] {len(question_files)} files × {chunks_per_file} chunks each")
                nodes, context, sources = await retrieve_per_file(
                    standalone, question_files, chunks_per_file
                )
                # For UML/PUML files: pin ALL nodes so no entity block is ever missing.
                retrieved_ids = {n.node_id for n in nodes}
                pinned_ctx: dict[str, list] = {}
                for fname in question_files:
                    if (fname.lower().split('.')[-1] in ('puml', 'plantuml', 'uml')
                            or _is_uml_image(fname)):
                        all_uml_nodes = get_nodes_for_files([fname])
                        extras = [n for n in all_uml_nodes if n.node_id not in retrieved_ids]
                        if extras:
                            pinned_ctx[fname] = extras
                            for n in extras:
                                retrieved_ids.add(n.node_id)
                            print(f"[pin] injected {len(extras)} UML/diagram chunks for {fname}")
                if pinned_ctx:
                    extra_parts = []
                    for fname, enodes in pinned_ctx.items():
                        extra_parts.append(
                            f"=== Document: {fname} (additional entities) ===\n"
                            + "\n\n".join(n.get_content() for n in enodes)
                        )
                    context = context + "\n\n" + "\n\n".join(extra_parts)
                    nodes = nodes + [n for extras in pinned_ctx.values() for n in extras]
            else:
                exhaustive = is_exhaustive_query(question)
                if exhaustive:
                    # For list/summary/extract queries, retrieve ALL chunks so nothing is missed.
                    # Cap at 200 to avoid overwhelming the LLM context window.
                    candidates_k = 200
                    print(f"[exhaustive] listing query detected — fetching up to {candidates_k} chunks")
                else:
                    candidates_k = SIMILARITY_TOP_K * 2 if reranker else SIMILARITY_TOP_K

                filters = MetadataFilters(
                    filters=[
                        MetadataFilter(key="file_name", value=fname, operator=FilterOperator.EQ)
                        for fname in question_files
                    ],
                    condition=FilterCondition.OR,
                )
                vector_retriever = index.as_retriever(
                    similarity_top_k=candidates_k, filters=filters
                )
                session_nodes = get_nodes_for_files(question_files)

                # HyDE + Multi-query: run both expansions in parallel, then retrieve.
                # • HyDE query  → vector search (semantic passage matching)
                # • Multi-query → additional vector searches (phrasing diversity)
                # • BM25        → keyword search on original standalone query
                expand_tasks = []
                if ENABLE_HYDE and not q.fast:
                    expand_tasks.append(hyde_expand(standalone, llm=aux_llm))
                if ENABLE_MULTI_QUERY and not q.fast:
                    expand_tasks.append(multi_query_expand(standalone, MULTI_QUERY_N, llm=aux_llm))

                expand_results = await asyncio.gather(*expand_tasks) if expand_tasks else []

                if _groq_daily_limit_hit and q.provider in ("cloud", "groq"):
                    yield f"data: {json.dumps({'type': 'error', 'message': 'Rate limit reached — Groq daily token quota exhausted. Please try again tomorrow or switch to a different model.'})}\n\n"
                    return

                idx = 0
                if ENABLE_HYDE and not q.fast:
                    hyde_query = expand_results[idx]; idx += 1
                    yield f"data: {json.dumps({'type': 'hypothesis', 'text': hyde_query})}\n\n"
                else:
                    hyde_query = standalone

                mq_queries = expand_results[idx] if (ENABLE_MULTI_QUERY and not q.fast and idx < len(expand_results)) else []

                if session_nodes:
                    bm25_retriever = BM25Retriever.from_defaults(
                        nodes=session_nodes, similarity_top_k=candidates_k
                    )
                    all_retrieve_coros = (
                        [vector_retriever.aretrieve(hyde_query)]
                        + [vector_retriever.aretrieve(q) for q in mq_queries]
                        + [bm25_retriever.aretrieve(standalone)]
                    )
                    all_results = await asyncio.gather(*all_retrieve_coros)
                    bm25_nodes    = all_results[-1]
                    all_vec_lists = all_results[:-1]

                    # Reciprocal rank fusion across all result lists
                    scores: dict[str, float] = {}
                    node_map: dict[str, object] = {}
                    for node_list in all_vec_lists:
                        for rank, n in enumerate(node_list):
                            scores[n.node_id] = scores.get(n.node_id, 0) + 1.0 / (rank + 60)
                            node_map[n.node_id] = n
                    for rank, n in enumerate(bm25_nodes):
                        scores[n.node_id] = scores.get(n.node_id, 0) + 1.0 / (rank + 60)
                        node_map[n.node_id] = n
                    nodes = sorted(node_map.values(), key=lambda n: scores[n.node_id], reverse=True)[:candidates_k]
                    print(f"[multi-query] merged {len(node_map)} unique chunks from {len(all_retrieve_coros)} queries")
                else:
                    nodes = await vector_retriever.aretrieve(hyde_query)

                if PARENT_CHILD_AVAILABLE:
                    retriever = AutoMergingRetriever(vector_retriever, storage_context, verbose=False)
                    try:
                        nodes = await retriever.aretrieve(hyde_query)
                    except Exception as _merge_err:
                        print(f"[warn] AutoMergingRetriever failed ({_merge_err}), using leaf nodes")

                if reranker and len(nodes) > 1:
                    _q, _k = standalone, SIMILARITY_TOP_K  # rerank against original query, not hypothesis

                    def _rerank(nodes_in, q=_q, k=_k):
                        pairs  = [(q, n.get_content()) for n in nodes_in]
                        scores = reranker.predict(pairs)
                        ranked = sorted(zip(nodes_in, scores), key=lambda x: x[1], reverse=True)
                        return [n for n, _ in ranked[:k]]

                    loop  = asyncio.get_event_loop()
                    nodes = await loop.run_in_executor(None, _rerank, nodes)
                    print(f"[rerank] kept {len(nodes)}/{candidates_k} chunks")

                # For PPTX files: always pin ALL overview chunks (the ones before
                # any "--- Slide" block). The overview is split across several child
                # chunks due to token limits, so we need all of them to guarantee
                # team names, slide count, etc. are always in the context.
                # For MLD files: always pin the overview chunk (schema_chunk == 'overview')
                # so exhaustive "list all entities" queries always get the full index,
                # even when top-K retrieval would otherwise omit it.
                retrieved_ids = {n.node_id for n in nodes}
                pinned = []
                for fname in question_files:
                    if fname.lower().endswith(".pptx"):
                        all_file_nodes = get_nodes_for_files([fname])
                        pptx_pin_limit = 5 if q.fast else len(all_file_nodes)
                        pptx_pinned = 0
                        for n in all_file_nodes:
                            content = n.get_content()
                            is_overview = not content.lstrip().startswith("--- Slide")
                            if is_overview and n.node_id not in retrieved_ids:
                                if pptx_pinned >= pptx_pin_limit:
                                    break
                                pinned.append(n)
                                retrieved_ids.add(n.node_id)
                                pptx_pinned += 1
                        print(f"[pin] injected {pptx_pinned} overview chunks for {fname}")
                    elif (fname.lower().split('.')[-1] in ('puml', 'plantuml', 'uml')
                          or _is_uml_image(fname)):
                        # UML/diagram files: pin ALL chunks — structured data,
                        # no partial retrieval allowed (later entity blocks must not be dropped).
                        all_file_nodes = get_nodes_for_files([fname])
                        uml_pinned = 0
                        for n in all_file_nodes:
                            if n.node_id not in retrieved_ids:
                                pinned.append(n)
                                retrieved_ids.add(n.node_id)
                                uml_pinned += 1
                        if uml_pinned:
                            print(f"[pin] injected {uml_pinned} UML/diagram chunks for {fname}")
                    else:
                        # Pin MLD overview nodes (any file that has schema_chunk='overview')
                        all_file_nodes = get_nodes_for_files([fname])
                        for n in all_file_nodes:
                            if n.metadata.get("schema_chunk") == "overview" and n.node_id not in retrieved_ids:
                                pinned.append(n)
                                retrieved_ids.add(n.node_id)
                                print(f"[pin] injected MLD overview chunk for {fname}")
                all_nodes_final = pinned + list(nodes)
                context = "\n\n".join(n.get_content() for n in all_nodes_final)
                sources = list({n.metadata.get("file_name", "unknown") for n in all_nodes_final})

            history_section = ""
            if history:
                history_lines = "\n".join(
                    f"User: {h.question}\nAssistant: {h.answer}" for h in history[-5:]
                )
                history_section = f"Conversation history:\n{history_lines}\n\n"

            if comparison_mode:
                src_list = ", ".join(f'"{f}"' for f in question_files)
                system_instruction = (
                    "You are a document assistant. You have been given excerpts from multiple documents, "
                    f"each clearly labeled with its filename: {src_list}.\n"
                    "Answer using ONLY the provided context.\n"
                    "CRITICAL RULES for multi-document answers:\n"
                    "- Always cite which document each piece of information comes from, "
                    "using the filename in parentheses e.g. (from mld.pdf) or (from diagram.puml).\n"
                    "- For UML / PlantUML files: extract and list EVERY entity/class with ALL its attributes "
                    "and ALL its relationships. Do not stop early — include every entity block in the context.\n"
                    "- When the question asks to explain or describe, cover both documents exhaustively — "
                    "do not omit any entity, class, or section that appears in the context.\n"
                    "- If a document lacks relevant information, say so explicitly.\n"
                    "- Format your response in clear Markdown: use numbered lists (`1.`, `2.`) or bullet lists (`-`) when enumerating items, use **bold** for key terms, and use headings (`##`) for long structured answers. Use a Markdown pipe table ONLY when every cell value fits on a single line with no nested lists or line breaks; if any cell would need multi-line content, use a bullet list instead."

                )
            else:
                system_instruction = (
                    "You are a document assistant. Answer using ONLY the provided context and conversation history.\n"
                    "Do not use any prior knowledge outside of these.\n"
                    "Important rules:\n"
                    "- Image files have already been analysed by a vision model. Their full description is in the context. "
                    "Do NOT ask the user to provide an image — all visual content is already available to you as text.\n"
                    "- Tables and row data are in the context as pipe-separated lines. Count or sum them directly from the text.\n"
                    "- When asked to list, enumerate, summarize, extract, or resume content, you MUST include "
                    "EVERY item that appears in the context. Do not stop early, do not skip entries, "
                    "do not use 'etc.', '...', or 'and more'. If you started listing, finish the list completely.\n"
                    "- If the answer cannot be found in the context, say exactly: "
                    "'I don't have enough information in the provided documents to answer this.'\n"
                    "- Do NOT add disclaimers, caveats, or meta-commentary about context limitations at the end of your answer. "
                    "Just answer directly.\n"
                    "- Format your response in clear Markdown: use numbered lists (`1.`, `2.`) or bullet lists (`-`) when enumerating items, use **bold** for key terms, and use headings (`##`) for long structured answers. Use a Markdown pipe table ONLY when every cell value fits on a single line with no nested lists or line breaks; if any cell would need multi-line content, use a bullet list instead."

                )

            exhaustive_reminder = (
                "IMPORTANT: This question asks for a complete list or summary. "
                "You MUST include every single item from the context above. "
                "Do not stop early. Do not use 'etc.' or '...'. List everything.\n\n"
            ) if exhaustive else ""

            final_prompt = (
                f"{system_instruction}"
                f"{history_section}"
                f"Context:\n---------------------\n{context}\n---------------------\n\n"
                f"{exhaustive_reminder}"
                f"Question: {question}\n\nAnswer:"
            )

            # 4. Stream character by character so the delay is uniform regardless
            #    of how many chars Ollama bundles into each chunk.
            STREAM_DELAY = float(os.getenv("STREAM_DELAY_MS", "20")) / 1000
            full_response = ""
            _t0 = asyncio.get_event_loop().time()
            async for chunk in await active_llm.astream_complete(final_prompt):
                if chunk.delta:
                    full_response += chunk.delta
                    for char in chunk.delta:
                        yield f"data: {json.dumps({'type': 'token', 'content': char})}\n\n"
                        await asyncio.sleep(STREAM_DELAY)

            # 5. Approximate token count for streaming (no raw field available per chunk)
            token_usage["completion"] += len(full_response.split())
            token_usage["requests"]   += 1

            _elapsed_ms = int((asyncio.get_event_loop().time() - _t0) * 1000)
            query_stats["total"]    += 1
            query_stats["total_ms"] += _elapsed_ms

            still_indexing = [f for f, s in indexing_status.items() if s == "indexing"]
            warning = (
                f"Still indexing: {', '.join(still_indexing)}. Results may be incomplete."
                if still_indexing else None
            )
            citations = build_citations(nodes)
            yield f"data: {json.dumps({'type': 'done', 'sources': sources, 'citations': citations, 'warning': warning, 'mode': 'comparison' if comparison_mode else 'standard'})}\n\n"

            if ENABLE_EVAL and not q.fast and full_response.strip():
                try:
                    eval_ctx = context[:3000]
                    eval_ans = full_response[:800]
                    print(f"[eval] ctx_len={len(context)} ans_len={len(full_response)}")
                    print(f"[eval] ctx_preview: {eval_ctx[:300]}")

                    faith_prompt = (
                        "You are evaluating a RAG system response.\n\n"
                        f"Retrieved context:\n{eval_ctx}\n\n"
                        f"AI answer:\n{eval_ans}\n\n"
                        "Faithfulness: are the answer's claims grounded in the context? "
                        "Count clearly hedged inferences (words like 'probably', 'likely', 'may', 'suggests') "
                        "as acceptable — they signal the model is not over-claiming. "
                        "Only penalise definite statements of fact that contradict or are absent from the context.\n"
                        "Reply with ONLY a decimal number from 0.0 (not faithful) to 1.0 (fully faithful)."
                    )
                    rel_prompt = (
                        "You are evaluating a RAG system response.\n\n"
                        f"User question: {question}\n\n"
                        f"AI answer:\n{eval_ans}\n\n"
                        "Answer relevance: does the answer directly and completely address the question?\n"
                        "Reply with ONLY a decimal number from 0.0 (irrelevant) to 1.0 (perfectly relevant)."
                    )

                    faith_result = await aux_llm.acomplete(faith_prompt)
                    print(f"[eval] faith_raw: {str(faith_result)[:80]}")
                    faith_score  = parse_eval_score(str(faith_result))

                    rel_result   = await aux_llm.acomplete(rel_prompt)
                    print(f"[eval] rel_raw: {str(rel_result)[:80]}")
                    rel_score    = parse_eval_score(str(rel_result))

                    faith_score = faith_score if faith_score is not None else 0.0
                    rel_score   = rel_score   if rel_score   is not None else 0.0
                    print(f"[eval] faithfulness={faith_score}  relevance={rel_score}")
                    yield f"data: {json.dumps({'type': 'eval', 'faithfulness': faith_score, 'answer_relevance': rel_score})}\n\n"
                except Exception as eval_err:
                    print(f"[eval] error: {eval_err}")

        except Exception as e:
            yield f"data: {json.dumps({'type': 'error', 'message': str(e)})}\n\n"

    return StreamingResponse(
        event_stream(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        }
    )

@app.post("/eval/quality")
async def run_quality_eval(
    provider:    str       = "cloud",
    groq_model:  str | None = None,
    current_user: UserModel = Depends(get_current_user),
):
    """
    Answer-quality evaluation: runs the full RAG pipeline for each question
    that has type=answer_quality + expected_answer, then scores with LLM:
      • faithfulness   – answer grounded in retrieved context?
      • relevance      – answer addresses the question?
      • correctness    – answer matches the expected answer?
    Returns per-question breakdown + aggregate scores.
    """
    dataset_path = os.path.join(os.path.dirname(__file__), "eval_dataset.json")
    if not os.path.exists(dataset_path):
        raise HTTPException(404, "eval_dataset.json not found")

    with open(dataset_path, encoding="utf-8") as f:
        raw = json.load(f)

    questions = [
        q for q in raw
        if isinstance(q, dict)
        and q.get("id")
        and not str(q.get("id", "")).startswith("_")
        and q.get("type") == "answer_quality"
        and q.get("expected_answer")
    ]

    if not questions:
        raise HTTPException(400, "No answer_quality questions with expected_answer found in eval_dataset.json")

    # ── LLMs ─────────────────────────────────────────────────────────────────
    try:
        main_llm = _get_llm(provider, groq_model=groq_model)
    except Exception:
        main_llm = Settings.llm
    try:
        scorer   = _get_llm("cloud", groq_model=GROQ_AUX_MODEL)   # always 8B for scoring
    except Exception:
        scorer   = Settings.llm

    # ── Retriever helpers (mirrors /eval pipeline) ────────────────────────────
    def _build_filters(source_files):
        return MetadataFilters(
            filters=[MetadataFilter(key="file_name", value=f, operator=FilterOperator.EQ)
                     for f in source_files],
            condition=FilterCondition.OR,
        )

    async def _retrieve(question, source_files, k=SIMILARITY_TOP_K):
        if not index:
            return []
        candidates = k * 2 if reranker else k
        kw = {"similarity_top_k": candidates}
        if source_files:
            kw["filters"] = _build_filters(source_files)
        vec_ret    = index.as_retriever(**kw)
        bm25_nodes = get_nodes_for_files(source_files) if source_files else []
        if bm25_nodes:
            bm25_ret  = BM25Retriever.from_defaults(nodes=bm25_nodes, similarity_top_k=candidates)
            retriever = QueryFusionRetriever([vec_ret, bm25_ret], similarity_top_k=candidates,
                                             num_queries=1, mode="reciprocal_rerank", use_async=False)
        else:
            retriever = vec_ret
        nodes = await retriever.aretrieve(question)
        if reranker and len(nodes) > 1:
            loop = asyncio.get_event_loop()
            def _rr(ns=list(nodes), q=question, top=k):
                pairs  = [(q, n.get_content()) for n in ns]
                scores = reranker.predict(pairs)
                ranked = sorted(zip(ns, scores), key=lambda x: x[1], reverse=True)
                return [n for n, _ in ranked[:top]]
            nodes = await loop.run_in_executor(None, _rr)
        return list(nodes)[:k]

    async def _score(prompt: str) -> float:
        try:
            result = await scorer.acomplete(prompt)
            s = parse_eval_score(str(result))
            return s if s is not None else 0.0
        except Exception:
            return 0.0

    # ── Main loop ─────────────────────────────────────────────────────────────
    results = []
    for q in questions:
        question_text = q["question"].strip()
        expected      = q["expected_answer"].strip()
        source_files  = q.get("source_files", [])

        nodes   = await _retrieve(question_text, source_files)
        context = "\n\n".join(n.get_content() for n in nodes)[:4000]

        # Generate answer
        try:
            gen_prompt = (
                f"Answer the following question based strictly on the provided context.\n\n"
                f"Context:\n{context}\n\n"
                f"Question: {question_text}\n\n"
                f"Answer concisely and directly."
            )
            gen_result  = await main_llm.acomplete(gen_prompt)
            gen_answer  = str(gen_result).strip()
        except Exception as e:
            gen_answer = f"[generation error: {e}]"

        # Score
        faith = await _score(
            f"Retrieved context:\n{context[:3000]}\n\n"
            f"AI answer:\n{gen_answer[:800]}\n\n"
            "Faithfulness: are the answer\'s claims grounded in the context? "
            "Penalise only definite statements absent from or contradicting the context. "
            "Reply with ONLY a decimal 0.0–1.0."
        )
        relevance = await _score(
            f"Question: {question_text}\n\n"
            f"AI answer:\n{gen_answer[:800]}\n\n"
            "Answer relevance: does the answer directly and completely address the question? "
            "Reply with ONLY a decimal 0.0–1.0."
        )
        correctness = await _score(
            f"Question: {question_text}\n\n"
            f"Expected answer:\n{expected}\n\n"
            f"AI answer:\n{gen_answer[:800]}\n\n"
            "Correctness: does the AI answer contain all key facts from the expected answer? "
            "Score 1.0 = all key facts present, 0.5 = partially correct, 0.0 = wrong or missing key facts. "
            "Reply with ONLY a decimal 0.0–1.0."
        )

        results.append({
            "id"           : q["id"],
            "question"     : question_text,
            "expected"     : expected,
            "generated"    : gen_answer,
            "source_files" : source_files,
            "faithfulness" : faith,
            "relevance"    : relevance,
            "correctness"  : correctness,
        })

    n = len(results)
    def _avg(key): return round(sum(r[key] for r in results) / n, 3) if n else 0.0

    return {
        "provider"      : provider,
        "model"         : groq_model or "local",
        "n_questions"   : n,
        "avg_faithfulness"  : _avg("faithfulness"),
        "avg_relevance"     : _avg("relevance"),
        "avg_correctness"   : _avg("correctness"),
        "per_question"  : results,
    }

@app.delete("/documents/all")
def clear_all(current_user: UserModel = Depends(_require_admin)):
    global index, storage_context
    chroma_client.delete_collection("rag_docs")
    global chroma_collection, vector_store
    chroma_collection = chroma_client.get_or_create_collection("rag_docs")
    vector_store = ChromaVectorStore(chroma_collection=chroma_collection)
    
    for f in os.listdir(UPLOAD_DIR):
        try:
            os.remove(f"{UPLOAD_DIR}/{f}")
        except:
            pass
            
    if PARENT_CHILD_AVAILABLE:
        docstore = SimpleDocumentStore()

        _ds_path = os.path.join(
        NODE_STORE_DIR,
        "docstore.json"
        )

        docstore.persist(persist_path=_ds_path)

        storage_context = StorageContext.from_defaults(
        vector_store=vector_store,
        docstore=docstore,
    )
    else:
        storage_context = StorageContext.from_defaults(
        vector_store=vector_store
    )
    index = None
    indexing_status.clear()
    return {"message": "All documents cleared"}

@app.delete("/documents/{filename}")
async def delete_document(filename: str,
                          current_user: UserModel = Depends(get_current_user)):
    if current_user.role != "admin" and file_owners.get(filename) != current_user.id:
        raise HTTPException(403, "Access denied")
    path = _safe_upload_path(filename)
    if os.path.exists(path):
        try:
            os.remove(path)
        except Exception as e:
            print(f"Failed to delete physical file {filename}: {e}")

    try:
        results = chroma_collection.get(
            where={"file_name": filename}
        )
        if results["ids"]:
            chroma_collection.delete(ids=results["ids"])
            print(f"Deleted {len(results['ids'])} chunks for {filename}")
    except Exception as e:
        print(f"ChromaDB delete error: {e}")

    global index, storage_context
    if chroma_collection.count() > 0:
        # Recreate storage context to clear in-memory node caches after deletion.
        if PARENT_CHILD_AVAILABLE:
            docstore = SimpleDocumentStore()

            _ds_path = os.path.join(
            NODE_STORE_DIR,
            "docstore.json"
            )

            docstore.persist(persist_path=_ds_path)

            storage_context = StorageContext.from_defaults(
            vector_store=vector_store,
            docstore=docstore,
            )
        else:
            storage_context = StorageContext.from_defaults(
            vector_store=vector_store
            )
            index = VectorStoreIndex.from_vector_store(
            vector_store,
            storage_context=storage_context
            )
    else:
        index = None

    indexing_status.pop(filename, None)
    file_owners.pop(filename, None)
    _save_file_owners()

    return {"deleted": filename}

@app.post("/reindex/{filename}")
async def reindex_document(filename: str,
                           current_user: UserModel = Depends(get_current_user)):
    """Force re-extraction and re-indexing of an already-uploaded file."""
    if current_user.role != "admin" and file_owners.get(filename) != current_user.id:
        raise HTTPException(403, "Access denied")

    path = os.path.join(UPLOAD_DIR, filename)
    if not os.path.exists(path):
        raise HTTPException(404, "File not found")

    try:
        results = chroma_collection.get(where={"file_name": filename})
        if results["ids"]:
            chroma_collection.delete(ids=results["ids"])
    except Exception as e:
        print(f"[reindex] ChromaDB cleanup error: {e}")

    # Clear hash so the upload endpoint won't skip re-indexing due to unchanged content.
    file_hashes.pop(filename, None)
    indexing_status[filename] = "indexing"

    loop = asyncio.get_event_loop()
    loop.run_in_executor(executor, add_document_to_index, path, filename, _active_provider)

    return {"status": "reindexing", "filename": filename}


@app.post("/eval")
async def run_eval(top_k: int = SIMILARITY_TOP_K):
    """
    Run the retrieval evaluation against eval_dataset.json.
    Returns per-question results, aggregate metrics, and a 3-way configuration
    comparison (Vector only / Hybrid / Hybrid + Reranker).
    """
    dataset_path = os.path.join(os.path.dirname(__file__), "eval_dataset.json")
    if not os.path.exists(dataset_path):
        raise HTTPException(404, "eval_dataset.json not found in rag-backend/")

    with open(dataset_path, encoding="utf-8") as f:
        raw = json.load(f)

    questions = [
        q for q in raw
        if isinstance(q, dict)
        and q.get("id")
        and not str(q.get("id", "")).startswith("_")
        and (q.get("answer_keywords") or q.get("source_files"))
    ]

    if not questions:
        raise HTTPException(400, "No evaluable questions in eval_dataset.json")

    def _is_relevant(node, keywords, sources):
        chunk_file = node.metadata.get("file_name", "")
        chunk_text = node.get_content().lower()
        if sources and chunk_file.lower() not in [s.lower() for s in sources]:
            return False
        if keywords:
            return any(kw.lower() in chunk_text for kw in keywords)
        return bool(sources)

    def _agg(hits, mrrs, n):
        hr = sum(hits) / n
        return {"hit_rate": round(hr, 3), "mrr": round(sum(mrrs) / n, 3)}

    def _build_filters(source_files):
        return MetadataFilters(
            filters=[MetadataFilter(key="file_name", value=f, operator=FilterOperator.EQ)
                     for f in source_files],
            condition=FilterCondition.OR,
        )


    async def _vec_only(question, source_files, k):
        """Pure vector similarity — no BM25, no reranker."""
        if not index: return []
        kwargs = {"similarity_top_k": k}
        if source_files:
            kwargs["filters"] = _build_filters(source_files)
        nodes = await index.as_retriever(**kwargs).aretrieve(question)
        return list(nodes)[:k]

    async def _hybrid(question, source_files, k):
        """Hybrid vector + BM25 fusion — no reranker."""
        if not index: return []
        if source_files:
            vec_ret    = index.as_retriever(similarity_top_k=k, filters=_build_filters(source_files))
            bm25_nodes = get_nodes_for_files(source_files)
        else:
            vec_ret    = index.as_retriever(similarity_top_k=k)
            bm25_nodes = []
        if bm25_nodes:
            bm25_ret  = BM25Retriever.from_defaults(nodes=bm25_nodes, similarity_top_k=k)
            retriever = QueryFusionRetriever([vec_ret, bm25_ret], similarity_top_k=k,
                                             num_queries=1, mode="reciprocal_rerank", use_async=False)
        else:
            retriever = vec_ret
        nodes = await retriever.aretrieve(question)
        return list(nodes)[:k]

    async def _full(question, source_files, k):
        """Hybrid + AutoMerging + cross-encoder reranker (production pipeline)."""
        if not index: return []
        candidates = k * 2 if reranker else k
        if source_files:
            vec_ret    = index.as_retriever(similarity_top_k=candidates, filters=_build_filters(source_files))
            bm25_nodes = get_nodes_for_files(source_files)
        else:
            vec_ret    = index.as_retriever(similarity_top_k=candidates)
            bm25_nodes = []
        if bm25_nodes:
            bm25_ret  = BM25Retriever.from_defaults(nodes=bm25_nodes, similarity_top_k=candidates)
            retriever = QueryFusionRetriever([vec_ret, bm25_ret], similarity_top_k=candidates,
                                             num_queries=1, mode="reciprocal_rerank", use_async=False)
        else:
            retriever = vec_ret
        nodes = await retriever.aretrieve(question)
        if reranker and len(nodes) > 1:
            loop = asyncio.get_event_loop()
            def _rr(ns=list(nodes), q=question, top=k):
                pairs  = [(q, n.get_content()) for n in ns]
                scores = reranker.predict(pairs)
                ranked = sorted(zip(ns, scores), key=lambda x: x[1], reverse=True)
                return [n for n, _ in ranked[:top]]
            nodes = await loop.run_in_executor(None, _rr)
        return list(nodes)[:k]


    per_question = []
    vec_hits, vec_mrrs   = [], []
    hyb_hits, hyb_mrrs   = [], []
    full_hits, full_mrrs = [], []

    for q in questions:
        keywords = q.get("answer_keywords", [])
        sources  = q.get("source_files", [])
        text     = q["question"].strip()

        nv, nh, nf = (
            await _vec_only(text, sources, top_k),
            await _hybrid (text, sources, top_k),
            await _full   (text, sources, top_k),
        )

        def _metrics(nodes):
            pos = [i+1 for i, n in enumerate(nodes) if _is_relevant(n, keywords, sources)]
            return (bool(pos),
                    round(len(pos) / top_k, 3),
                    round(1.0 / min(pos), 3) if pos else 0.0,
                    min(pos) if pos else None)

        vh, vp, vm, _  = _metrics(nv)
        hh, hp, hm, _  = _metrics(nh)
        fh, fp, fm, fr = _metrics(nf)

        vec_hits.append(vh);  vec_mrrs.append(vm)
        hyb_hits.append(hh);  hyb_mrrs.append(hm)
        full_hits.append(fh); full_mrrs.append(fm)

        retrieved = [
            {
                "file": n.metadata.get("file_name", "?"),
                "page": str(n.metadata.get("page_label",
                            n.metadata.get("page_number", "?"))),
                "hit" : _is_relevant(n, keywords, sources),
            }
            for n in nf
        ]

        per_question.append({
            "id"          : q["id"],
            "question"    : text,
            "source_files": sources,
            "hit"         : fh,
            "precision"   : fp,
            "mrr"         : fm,
            "first_rank"  : fr,
            "retrieved"   : retrieved,
        })

    n = len(per_question)
    full_hr  = sum(full_hits) / n
    full_prec = sum(r["precision"] for r in per_question) / n
    full_mrr  = sum(r["mrr"]       for r in per_question) / n

    return {
        "top_k"          : top_k,
        "n_questions"    : n,
        "hit_rate"       : round(full_hr,   3),
        "precision"      : round(full_prec, 3),
        "recall"         : round(full_hr,   3),
        "mrr"            : round(full_mrr,  3),
        "configurations" : [
            {"name": "Vector only",        **_agg(vec_hits,  vec_mrrs,  n)},
            {"name": "Hybrid",             **_agg(hyb_hits,  hyb_mrrs,  n)},
            {"name": "Hybrid + Reranker",  **_agg(full_hits, full_mrrs, n)},
        ],
        "per_question"   : per_question,
    }



